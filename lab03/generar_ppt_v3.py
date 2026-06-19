"""
Presentación Caso 3 - LegendaryClass
Diseño basado exactamente en el estilo visual de la web:
  - slate-900 / purple-900 / slate-800 / slate-700
  - Tipografía limpia, menos elementos por slide, más espacio
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN

# ── Paleta exacta de la web ──────────────────────────────
SLATE_900 = RGB(15,  23,  42)   # bg principal
SLATE_800 = RGB(30,  41,  59)   # cards
SLATE_700 = RGB(51,  65,  85)   # borders / separadores
SLATE_400 = RGB(148,163,184)    # texto muted
SLATE_300 = RGB(203,213,225)    # texto secundario

PURPLE_950= RGB(59,  7, 100)    # purple-950 (via gradient)
PURPLE_900= RGB(74,  29,150)    # purple-900
PURPLE_600= RGB(147, 51,234)    # acento principal
PURPLE_500= RGB(168, 85,247)    # acento hover
PURPLE_400= RGB(192,132,252)    # acento suave
PURPLE_300= RGB(216,180,254)    # texto acento

INDIGO_600= RGB(79,  70,229)    # from-purple-600 to-indigo-600

WHITE     = RGB(255,255,255)
YELLOW_400= RGB(250,204, 21)    # XP / niveles (gold)
YELLOW_500= RGB(234,179,  8)
GREEN_400 = RGB(74, 222,128)
GREEN_500 = RGB(34, 197, 94)
RED_400   = RGB(248,113,113)

# ── Datos ────────────────────────────────────────────────
INSTITUCION = "TECSUP"
CARRERA     = "Ingeniería de Software  ·  V Ciclo"
CURSO       = "Diseño de Proyectos de Innovación"
DOCENTE     = "Turkowsky Vizcarra, Luisa"
ALUMNO      = "Aguirre, Juan"
ANIO        = "2026 - I"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────
def bg(sl):
    """Fondo degradado simulado slate-900 → purple-950 → slate-900"""
    _rect(sl, 0, 0, 13.33, 7.5, fill=SLATE_900)
    # Orbe izquierdo
    _rect(sl, -1.5, -1.5, 5, 5, fill=PURPLE_950)   # simula blur con opacidad baja
    # Orbe derecho
    _rect(sl, 10, 4, 5, 5, fill=PURPLE_950)

def _rect(sl, l, t, w, h, fill=None, line_color=None, line_w=12700):
    sp = sl.shapes.add_shape(1,
        Inches(max(l,-2)), Inches(max(t,-2)),
        Inches(min(w,16)),  Inches(min(h,12)))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Emu(line_w)
    else:
        sp.line.fill.background()
    return sp

def card(sl, l, t, w, h, border=True):
    """Card estilo slate-800 con borde slate-700"""
    sp = _rect(sl, l, t, w, h, fill=SLATE_800,
               line_color=SLATE_700 if border else None, line_w=9525)
    return sp

def txt(sl, text, l, t, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def pill(sl, text, l, t, size=10, bg_color=PURPLE_600, fg=WHITE, bold=False):
    """Badge / pill pequeño"""
    w = len(text) * 0.085 + 0.25
    _rect(sl, l, t, w, 0.32, fill=bg_color)
    txt(sl, text, l+0.05, t+0.02, w-0.05, 0.28,
        size=size, bold=bold, color=fg, align=PP_ALIGN.CENTER)
    return w

def accent_bar(sl, l, t, w, color=PURPLE_600, h=0.04):
    _rect(sl, l, t, w, h, fill=color)

def footer(sl):
    accent_bar(sl, 0, 7.36, 13.33, color=SLATE_700, h=0.04)
    txt(sl, f"LegendaryClass  ·  {CURSO}  ·  {INSTITUCION}  ·  {ANIO}",
        0.4, 7.4, 12.5, 0.22, size=8, color=SLATE_400)

def slide_header(sl, title, subtitle=""):
    """Header minimalista: solo texto, sin barra de color"""
    txt(sl, title, 0.55, 0.22, 12, 0.6,
        size=24, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.55, 0.82, 12, 0.38,
            size=12, color=PURPLE_300, italic=False)
    accent_bar(sl, 0.55, 1.25, 12.23, color=SLATE_700)
    footer(sl)


# ═══════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# Panel central limpio
card(sl, 2.5, 1.2, 8.33, 5.1)

# ⚔️ ícono
txt(sl, "⚔️", 5.6, 1.5, 2.2, 0.9, size=42, align=PP_ALIGN.CENTER)

# Título igual que la web: blanco bold
txt(sl, "LegendaryClass", 2.7, 2.45, 7.9, 0.95,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtítulo: purple-300
txt(sl, "Sistema SaaS de Gamificación Educativa",
    2.7, 3.42, 7.9, 0.5,
    size=16, color=PURPLE_300, align=PP_ALIGN.CENTER)

accent_bar(sl, 3.5, 4.0, 6.33, color=SLATE_700)

# Datos académicos
txt(sl, f"{CURSO}  ·  Caso 3 — Parte Práctica",
    2.7, 4.15, 7.9, 0.42,
    size=12, color=SLATE_300, align=PP_ALIGN.CENTER)
txt(sl, f"Docente: {DOCENTE}   ·   Alumno: {ALUMNO}",
    2.7, 4.58, 7.9, 0.38,
    size=11, color=SLATE_400, align=PP_ALIGN.CENTER)
txt(sl, f"{INSTITUCION}  ·  {CARRERA}  ·  {ANIO}",
    2.7, 4.98, 7.9, 0.35,
    size=11, color=PURPLE_400, align=PP_ALIGN.CENTER)

footer(sl)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — DESCRIPCIÓN DEL PRODUCTO
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Descripción del Producto", "LegendaryClass — plataforma SaaS de gamificación para docentes")

# Descripción principal
txt(sl, "Transforma el aula en un RPG educativo. Los estudiantes eligen clase de "
        "personaje, acumulan XP, suben de nivel y canjean recompensas. El docente "
        "gestiona comportamientos, misiones y reportes.",
    0.55, 1.42, 12.23, 0.85,
    size=13, color=SLATE_300, wrap=True)

# 4 cards de características
features = [
    ("⚔️", "5 Clases RPG",     "Mago · Guerrero · Ninja\nArquero · Lanzador\n+20 % XP por acción"),
    ("📈", "Motor de XP",       "Niveles, racha diaria\nlogros automáticos\nmisiones configurables"),
    ("👩‍🏫","Panel Docente",     "Comportamientos\nrecompensas y\nreportes por aula"),
    ("🏛️","Panel Director",    "Estadísticas globales\nde la institución\nen tiempo real"),
]
for i, (ico, title, desc) in enumerate(features):
    x = 0.55 + i * 3.2
    card(sl, x, 2.45, 3.0, 2.85)
    txt(sl, ico,   x+0.15, 2.55, 0.7,  0.55, size=26)
    txt(sl, title, x+0.15, 3.15, 2.7,  0.45, size=13, bold=True, color=WHITE)
    txt(sl, desc,  x+0.15, 3.6,  2.7,  1.55, size=10.5, color=SLATE_400, wrap=True)

# Stack + precio
accent_bar(sl, 0.55, 5.48, 12.23, color=SLATE_700)
txt(sl, "Stack:  Angular 18  ·  NestJS  ·  PostgreSQL  ·  Prisma  ·  TailwindCSS",
    0.55, 5.57, 8, 0.38, size=11.5, bold=True, color=PURPLE_300)
txt(sl, "Precio:  S/. 30 – 85 / mes por docente   ·   Sin costo para estudiantes",
    0.55, 5.97, 10, 0.35, size=11, color=SLATE_400)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — METODOLOGÍA
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Metodología", "Encuesta a 100 docentes de Lima Metropolitana · feb–mar 2026")

# Info encuesta
card(sl, 0.55, 1.42, 5.8, 1.75)
txt(sl, "Encuesta de campo", 0.75, 1.52, 5.4, 0.42, size=13, bold=True, color=PURPLE_300)
for j, item in enumerate([
    "100 docentes · Lima Metropolitana",
    "Google Forms · 20 preguntas",
    "Periodo: febrero – marzo 2026",
]):
    txt(sl, "· " + item, 0.75, 1.98+j*0.37, 5.4, 0.34, size=12, color=SLATE_300)

# Preguntas clave
card(sl, 6.65, 1.42, 6.13, 1.75)
txt(sl, "Preguntas clave", 6.85, 1.52, 5.7, 0.42, size=13, bold=True, color=PURPLE_300)
qs = [
    "Q4  —  Situación laboral            → Mercado Disponible",
    "Q6  —  Baja motivación             → Mercado Efectivo",
    "Q13 —  Disposición gamificación  → Mercado Objetivo",
    "Q17 —  Disposición de pago        → Precio ponderado",
]
for j, q in enumerate(qs):
    txt(sl, q, 6.85, 1.98+j*0.37, 5.8, 0.34, size=10.5, color=SLATE_300)

# Título cascada
txt(sl, "Modelo de segmentación en cascada",
    0.55, 3.35, 12, 0.42, size=14, bold=True, color=WHITE)

# 4 pasos cascada — clean
cascade = [
    ("Potencial",  "185,000", "MINEDU 2024",             PURPLE_600),
    ("Disponible", "177,600", "× 96 %  (Q4)",            PURPLE_500),
    ("Efectivo",   "103,008", "× 58 %  (Q6)",            PURPLE_400),
    ("Objetivo",    "82,406", "× 80 %  (Q13)",           YELLOW_400),
]
for i, (lbl, val, crit, col) in enumerate(cascade):
    x = 0.55 + i * 3.2
    card(sl, x, 3.85, 3.0, 2.2)
    accent_bar(sl, x, 3.85, 3.0, color=col)
    txt(sl, lbl,  x+0.15, 3.97, 2.7, 0.38, size=11, bold=True, color=col)
    txt(sl, val,  x+0.15, 4.38, 2.7, 0.65, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, crit, x+0.15, 5.05, 2.7, 0.35, size=10.5, color=SLATE_400, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "›", x+3.0, 4.55, 0.2, 0.5, size=20, bold=True,
            color=SLATE_700, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS DE LA ENCUESTA
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Resultados de la Encuesta", "Preguntas clave para la segmentación de mercado  ·  n = 100")

# Q6 — barras horizontales (izquierda)
txt(sl, "Q6  —  Frecuencia de baja motivación estudiantil",
    0.55, 1.42, 7.0, 0.42, size=13, bold=True, color=WHITE)

q6_data = [
    ("Siempre",        22, PURPLE_600),
    ("Frecuentemente", 36, PURPLE_500),
    ("A veces",        30, SLATE_700),
    ("Raramente",       8, SLATE_700),
    ("Nunca",           4, SLATE_700),
]
for j, (lbl, val, col) in enumerate(q6_data):
    y = 1.95 + j * 0.73
    # fondo barra
    _rect(sl, 0.55, y, 6.5, 0.55, fill=SLATE_800)
    # barra
    bar_w = (val / 36) * 5.8
    _rect(sl, 0.55, y, bar_w, 0.55, fill=col)
    # etiqueta
    txt(sl, lbl,      0.68, y+0.12, 2.0, 0.3, size=11,
        bold=(j<2), color=WHITE if j<2 else SLATE_400)
    txt(sl, f"{val} %", 0.55+bar_w+0.1, y+0.12, 0.7, 0.3,
        size=12, bold=(j<2), color=col if col!=SLATE_700 else SLATE_400)
    if j < 2:
        txt(sl, "← filtro", 0.55+bar_w+0.82, y+0.14, 1.5, 0.28,
            size=9, color=PURPLE_400, italic=True)

# Resultado = 58%
card(sl, 0.55, 5.6, 6.5, 0.85)
txt(sl, "Siempre + Frecuentemente  =  58 %  →  Mercado Efectivo",
    0.75, 5.73, 6.1, 0.45, size=12, bold=True, color=PURPLE_300)

# Panel derecho — resultados clave Q4, Q13, Q17
txt(sl, "Hallazgos por pregunta", 7.7, 1.42, 5.2, 0.42,
    size=13, bold=True, color=WHITE)

hallazgos = [
    ("Q4",  "96 %",  "Trabajan en institución educativa",    PURPLE_600),
    ("Q6",  "58 %",  "Baja motivación frecuente/siempre",    PURPLE_500),
    ("Q13", "80 %",  "Implementarían gamificación",           GREEN_400),
    ("Q17", "61/100","Con intención definitiva de pago",      YELLOW_400),
]
for j, (q, val, desc, col) in enumerate(hallazgos):
    y = 1.95 + j * 1.12
    card(sl, 7.7, y, 5.2, 0.95)
    accent_bar(sl, 7.7, y, 5.2, color=col)
    txt(sl, q,   7.85, y+0.12, 0.6, 0.4, size=10, bold=True, color=col)
    txt(sl, val, 8.5,  y+0.08, 1.8, 0.52, size=22, bold=True, color=WHITE)
    txt(sl, desc, 10.3, y+0.18, 2.4, 0.5, size=10, color=SLATE_400, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — ANÁLISIS DE PRECIO
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Análisis de Precio", "Q17 — Disposición de pago  ·  61 respondentes con intención definitiva de pago")

# 3 cards de precio
price_cards = [
    ("S/. 20 – 40 / mes",  "S/. 30",  "S/. 360 / año",  35, PURPLE_600),
    ("S/. 41 – 70 / mes",  "S/. 55",  "S/. 660 / año",  18, PURPLE_500),
    ("S/. 71 – 100 / mes", "S/. 85", "S/. 1,020 / año",  8, PURPLE_400),
]
for i, (rng, mo, yr, qty, col) in enumerate(price_cards):
    x = 0.55 + i * 4.1
    card(sl, x, 1.42, 3.85, 3.5)
    accent_bar(sl, x, 1.42, 3.85, color=col)
    txt(sl, rng, x+0.15, 1.54, 3.55, 0.38, size=11, color=SLATE_400, align=PP_ALIGN.CENTER)
    txt(sl, mo,  x+0.15, 1.96, 3.55, 0.85, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    accent_bar(sl, x+0.3, 2.88, 3.25, color=SLATE_700)
    txt(sl, yr,  x+0.15, 2.98, 3.55, 0.52, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, f"{qty} docentes", x+0.15, 3.55, 3.55, 0.38,
        size=13, color=SLATE_400, align=PP_ALIGN.CENTER)

# Resultado precio ponderado — card destacado
card(sl, 0.55, 5.1, 12.23, 1.45)
accent_bar(sl, 0.55, 5.1, 12.23, color=YELLOW_500)

txt(sl, "Precio Ponderado Anual", 0.8, 5.2, 5.5, 0.42,
    size=12, color=YELLOW_400)
txt(sl, "S/. 535.08 / docente / año   ≈   S/. 44.59 / mes",
    0.8, 5.62, 9.5, 0.7, size=28, bold=True, color=WHITE)
txt(sl, "Cálculo:  ( 35 × S/.360  +  18 × S/.660  +  8 × S/.1,020 )  ÷  61  =  S/.32,640 ÷ 61",
    0.8, 6.32, 11.5, 0.35, size=10.5, color=SLATE_400, italic=True)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — ESCENARIOS
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Análisis de la Demanda — Escenarios", "Año 2027  ·  Mercado Objetivo: 82,406 docentes  ·  Precio: S/. 535.08 / año")

esc = [
    ("Optimista",  "80 %", "65,925",  "S/. 35,274,015",  GREEN_400),
    ("Moderado",   "50 %", "41,203",  "S/. 22,046,259",  YELLOW_400),
    ("Pesimista",  "30 %", "24,722",  "S/. 13,227,786",  RED_400),
]
for i, (name, pct, subs, rev, col) in enumerate(esc):
    x = 0.55 + i * 4.15
    card(sl, x, 1.42, 3.9, 5.3)
    accent_bar(sl, x, 1.42, 3.9, color=col)
    # Nombre
    txt(sl, name,  x+0.2, 1.55, 3.5, 0.45, size=16, bold=True, color=col)
    # Porcentaje grande
    txt(sl, pct,   x+0.2, 2.05, 3.5, 1.1,  size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "penetración", x+0.2, 3.15, 3.5, 0.38, size=10, color=SLATE_400, align=PP_ALIGN.CENTER)
    accent_bar(sl, x+0.3, 3.62, 3.3, color=SLATE_700)
    txt(sl, "Suscripciones", x+0.2, 3.74, 3.5, 0.35, size=10, color=SLATE_400, align=PP_ALIGN.CENTER)
    txt(sl, subs,  x+0.2, 4.1,  3.5, 0.62, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    accent_bar(sl, x+0.3, 4.8, 3.3, color=SLATE_700)
    txt(sl, "Ingresos anuales", x+0.2, 4.92, 3.5, 0.35, size=10, color=SLATE_400, align=PP_ALIGN.CENTER)
    txt(sl, rev,   x+0.2, 5.3,  3.5, 0.52, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(sl, "Demanda = Mercado Objetivo × % Penetración   ·   Ingreso = Demanda × S/. 535.08",
    0.55, 6.9, 12.23, 0.32, size=10, color=SLATE_400, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 7 — PROYECCIÓN 5 AÑOS
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Proyección 2027–2031", "Escenario Moderado  ·  Crecimiento anual: +15 %  ·  HolonIQ EdTech LATAM 2024")

# KPIs
kpis = [
    ("41,203 → 72,063", "Suscripciones  2027 → 2031",  PURPLE_300),
    ("+75 %",           "Crecimiento acumulado",         GREEN_400),
    ("S/. 148.7M",      "Ingresos acumulados 2027–31",   YELLOW_400),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.55 + i * 4.1
    card(sl, x, 1.42, 3.85, 1.25)
    accent_bar(sl, x, 1.42, 3.85, color=col)
    txt(sl, val, x+0.2, 1.54, 3.45, 0.58, size=20, bold=True, color=col)
    txt(sl, lbl, x+0.2, 2.1,  3.45, 0.42, size=10, color=SLATE_400)

# Gráfico de barras — clean
PROJ_YEARS = [2027, 2028, 2029, 2030, 2031]
PROJ_SUBS  = [41203, 47383, 54490, 62664, 72063]
PROJ_ING   = [22046259, 25353198, 29156177, 33529604, 38559045]
bar_colors = [PURPLE_600, PURPLE_500, PURPLE_400, PURPLE_300, YELLOW_400]

max_ing    = max(PROJ_ING)
chart_b    = 6.6
chart_h    = 3.5
bar_w      = 1.75
gap        = 0.45
start_x    = 0.8

# Línea base
accent_bar(sl, start_x-0.1, chart_b, 12, color=SLATE_700)

for i, (yr, sub, ing, col) in enumerate(zip(PROJ_YEARS, PROJ_SUBS, PROJ_ING, bar_colors)):
    x     = start_x + i * (bar_w + gap)
    bh    = (ing / max_ing) * chart_h
    y_bar = chart_b - bh
    _rect(sl, x, y_bar, bar_w, bh, fill=col)
    # valor sobre barra
    ing_m = ing / 1_000_000
    txt(sl, f"S/.{ing_m:.1f}M", x-0.1, y_bar-0.45, bar_w+0.2, 0.38,
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    # año y suscripciones bajo la barra
    txt(sl, str(yr),          x, chart_b+0.06, bar_w, 0.35,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, f"{sub:,}",       x, chart_b+0.42, bar_w, 0.3,
        size=9.5, color=SLATE_400, align=PP_ALIGN.CENTER)

# Nota
txt(sl, "Fuente de tasa de crecimiento: HolonIQ EdTech LATAM Report 2024  ·  CAGR 15 % para SaaS educativo en Perú/LATAM",
    0.55, 7.1, 12.23, 0.28, size=8.5, color=SLATE_400, italic=True)

# ═══════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSIONES
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_header(sl, "Conclusiones", "Síntesis del análisis de demanda — LegendaryClass")

concl = [
    (PURPLE_300, "82,406 docentes",   "conforman el mercado objetivo en Lima Metropolitana."),
    (YELLOW_400, "S/. 535 / año",     "precio ponderado validado por disposición de pago real (Q17)."),
    (GREEN_400,  "S/. 22M en 2027",   "ingresos proyectados en escenario moderado (50 % penetración)."),
    (PURPLE_400, "S/. 38.5M en 2031", "con crecimiento anual del 15 % (CAGR EdTech LATAM)."),
    (PURPLE_300, "58 % de docentes",  "reportan baja motivación frecuente — necesidad validada por encuesta."),
]
for i, (col, val, rest) in enumerate(concl):
    y = 1.42 + i * 1.05
    card(sl, 0.55, y, 12.23, 0.92)
    accent_bar(sl, 0.55, y, 0.08, color=col, h=0.92)   # borde izquierdo
    # Número
    txt(sl, str(i+1), 0.75, y+0.2, 0.48, 0.5,
        size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Valor destacado
    txt(sl, val,  1.35, y+0.18, 3.2, 0.55, size=17, bold=True, color=col)
    # Descripción
    txt(sl, rest, 4.6,  y+0.22, 7.9, 0.5,  size=13, color=SLATE_300)

# ═══════════════════════════════════════════════════════
# SLIDE 9 — CIERRE
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

card(sl, 3.2, 2.0, 6.93, 3.5)
accent_bar(sl, 3.2, 2.0, 6.93, color=PURPLE_600)

txt(sl, "⚔️", 5.9, 2.15, 1.5, 0.7, size=32, align=PP_ALIGN.CENTER)
txt(sl, "LegendaryClass", 3.4, 2.9, 6.53, 0.85,
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Sistema SaaS de Gamificación Educativa",
    3.4, 3.78, 6.53, 0.45, size=13, color=PURPLE_300, align=PP_ALIGN.CENTER)
accent_bar(sl, 4.0, 4.32, 5.33, color=SLATE_700)
txt(sl, "¡Gracias!",
    3.4, 4.45, 6.53, 0.72, size=30, bold=True, color=YELLOW_400, align=PP_ALIGN.CENTER)

txt(sl, f"{ALUMNO}  ·  {CURSO}  ·  {ANIO}",
    3.4, 5.28, 6.53, 0.38, size=10.5, color=SLATE_400, align=PP_ALIGN.CENTER)

footer(sl)

# ── Guardar ──────────────────────────────────────────────
OUT = "C:/Proyectos/LegendaryClass/lab03/Presentacion_Caso3_LegendaryClass_v3.pptx"
prs.save(OUT)
print("[OK] PPT guardada:", OUT)
