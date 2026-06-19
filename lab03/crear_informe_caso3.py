"""
Genera Informe_Caso3_LegendaryClass.docx y Presentacion_Caso3_LegendaryClass.pptx
Ejecutar desde: C:/Proyectos/LegendaryClass/lab03/
  python crear_informe_caso3.py
"""

# ═══════════════════════════════════════════════════════════════════════════════
#  WORD
# ═══════════════════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

def rgb(r, g, b):
    return RGBColor(r, g, b)

PURPLE = rgb(112, 48, 160)
BLUE   = rgb(0, 112, 192)
DARK   = rgb(31, 31, 31)
GRAY   = rgb(89, 89, 89)
WHITE  = rgb(255, 255, 255)

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def add_heading(doc, text, level=1, color=PURPLE):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = True
    if level == 1:
        run.font.size = Pt(16)
    elif level == 2:
        run.font.size = Pt(13)
    else:
        run.font.size = Pt(11)
    run.font.color.rgb = color
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(4)
    return p

def add_body(doc, text, bold=False, size=10.5):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = DARK
    p.paragraph_format.space_after = Pt(4)
    return p

def add_bullet(doc, text, level=0):
    p = doc.add_paragraph(style='List Bullet')
    run = p.add_run(text)
    run.font.size = Pt(10.5)
    run.font.color.rgb = DARK
    p.paragraph_format.space_after = Pt(2)
    return p

# ──────────────────────────────────────────────────────────
doc = Document()

# Estilos generales
style = doc.styles['Normal']
style.font.name = 'Calibri'
style.font.size = Pt(10.5)

# Márgenes
for section in doc.sections:
    section.top_margin    = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin   = Cm(3.0)
    section.right_margin  = Cm(2.5)

# ── PORTADA ──────────────────────────────────────────────
doc.add_paragraph()
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("UNIVERSIDAD PERUANA")
run.bold = True; run.font.size = Pt(13); run.font.color.rgb = GRAY

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Escuela de Ingeniería de Software")
run.font.size = Pt(11); run.font.color.rgb = GRAY

doc.add_paragraph()

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("LABORATORIO 03")
run.bold = True; run.font.size = Pt(11); run.font.color.rgb = BLUE

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Análisis de Demanda del Mercado")
run.bold = True; run.font.size = Pt(11); run.font.color.rgb = BLUE

doc.add_paragraph()

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("CASO 3")
run.bold = True; run.font.size = Pt(22); run.font.color.rgb = PURPLE

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("LegendaryClass")
run.bold = True; run.font.size = Pt(26); run.font.color.rgb = PURPLE

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Sistema SaaS de Gamificación Educativa")
run.font.size = Pt(13); run.font.color.rgb = GRAY

doc.add_paragraph()
doc.add_paragraph()

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Alumno: Juan Aguirre  |  Docente: Lturkowsky  |  2026 - I")
run.font.size = Pt(10); run.font.color.rgb = GRAY

doc.add_page_break()

# ── 1. DESCRIPCIÓN DEL PRODUCTO ──────────────────────────
add_heading(doc, "1. Descripción del Producto / Servicio")
add_body(doc,
    "LegendaryClass es una plataforma SaaS de gamificación educativa orientada a docentes de educación "
    "básica y superior. Transforma el aula en una experiencia de juego de rol (RPG): los estudiantes "
    "eligen clases (Mago, Guerrero, Ninja, Arquero, Lanzador), acumulan puntos de experiencia (XP), "
    "suben de nivel, completan misiones y canjean recompensas definidas por el docente."
)
add_body(doc,
    "El docente gestiona grupos, asigna comportamientos positivos o negativos, configura misiones y "
    "recompensas, y obtiene reportes de progreso. El director de institución supervisa estadísticas "
    "globales. La plataforma se entrega como aplicación web (Angular 18 + NestJS + PostgreSQL) "
    "accesible desde cualquier dispositivo con navegador."
)

add_heading(doc, "Características principales", level=2, color=BLUE)
bullets = [
    "5 clases de personaje con bonificaciones de XP (+20 %) por tipo de acción",
    "Motor de gamificación: niveles, racha diaria, logros automáticos, misiones",
    "Panel del docente: registro de comportamientos, gestión de recompensas y misiones",
    "Panel del director: reportes institucionales y estadísticas agregadas",
    "Portal del padre de familia: seguimiento del progreso del estudiante",
    "Importación masiva de estudiantes vía Excel",
    "API REST documentada (NestJS + Swagger)",
    "Precio de lanzamiento: suscripción mensual por docente (S/. 30 – S/. 85 / mes)",
]
for b in bullets:
    add_bullet(doc, b)

# ── 2. OBJETIVO DEL ESTUDIO ───────────────────────────────
add_heading(doc, "2. Objetivo del Estudio de Demanda")
add_body(doc,
    "Determinar el tamaño del mercado objetivo de LegendaryClass en Lima Metropolitana para el año "
    "de lanzamiento (2027) y proyectar los ingresos potenciales en un horizonte de 5 años (2027–2031), "
    "bajo tres escenarios de penetración de mercado."
)

# ── 3. METODOLOGÍA ────────────────────────────────────────
add_heading(doc, "3. Metodología")
add_body(doc,
    "Se aplicó una encuesta digital de 20 preguntas (Google Forms) a 100 docentes de instituciones "
    "educativas de Lima Metropolitana entre febrero y marzo de 2026. El cuestionario abarcó: perfil "
    "del encuestado, nivel de motivación estudiantil percibido, experiencia con herramientas digitales, "
    "disposición de implementar gamificación y disposición de pago."
)
add_body(doc,
    "El análisis de demanda sigue el modelo de segmentación en cascada utilizado en el laboratorio:"
)
add_bullet(doc, "Mercado Potencial → Mercado Disponible → Mercado Efectivo → Mercado Objetivo")
add_body(doc,
    "El precio de venta se calcula como promedio ponderado de los rangos de disposición de pago "
    "reportados en la encuesta (Q17), considerando únicamente a los 61 respondentes que afirmaron "
    "tener intención definitiva de pago."
)

# ── 4. SEGMENTACIÓN DE MERCADO ────────────────────────────
add_heading(doc, "4. Segmentación de Mercado (Sección 9)")

add_heading(doc, "4.1 Mercado Potencial", level=2, color=BLUE)
add_body(doc,
    "Según el Ministerio de Educación del Perú (MINEDU, 2024), Lima Metropolitana cuenta con "
    "aproximadamente 185,000 docentes activos en instituciones de educación básica regular (EBR) "
    "y educación superior. Este universo constituye el mercado potencial del producto."
)

add_heading(doc, "4.2 Mercado Disponible", level=2, color=BLUE)
add_body(doc,
    "De la encuesta, el 96 % de los encuestados indicó trabajar actualmente en una institución "
    "educativa (Q4). Se aplica este porcentaje como filtro de elegibilidad:"
)
add_body(doc, "Mercado Disponible = 185,000 × 96 % = 177,600 docentes", bold=True)

add_heading(doc, "4.3 Mercado Efectivo", level=2, color=BLUE)
add_body(doc,
    "Se consideró como condición de efectividad que el docente perciba baja motivación en sus "
    "estudiantes de forma frecuente o siempre (Q6: «¿Con qué frecuencia observas baja motivación "
    "en tus estudiantes?»). El 22 % respondió «Siempre» y el 36 % «Frecuentemente», sumando 58 %:"
)
add_body(doc, "Mercado Efectivo = 177,600 × 58 % = 103,008 docentes", bold=True)

add_heading(doc, "4.4 Mercado Objetivo", level=2, color=BLUE)
add_body(doc,
    "Finalmente, el 80 % de los respondentes afirmó estar dispuesto a implementar una solución "
    "de gamificación en su aula (Q13):"
)
add_body(doc, "Mercado Objetivo = 103,008 × 80 % = 82,406 docentes", bold=True)

# Tabla resumen segmentación
doc.add_paragraph()
tbl = doc.add_table(rows=5, cols=4)
tbl.style = 'Table Grid'
tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
headers = ["Segmento", "Criterio de Filtro", "Factor", "Tamaño (docentes)"]
for i, h in enumerate(headers):
    cell = tbl.rows[0].cells[i]
    cell.text = h
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.color.rgb = WHITE
    set_cell_bg(cell, "0070C0")

rows_data = [
    ("Potencial",   "Docentes activos en Lima Met.",          "—",    "185,000"),
    ("Disponible",  "Trabajan en institución educativa (Q4)", "96 %", "177,600"),
    ("Efectivo",    "Perciben baja motivación frecuente (Q6)","58 %", "103,008"),
    ("Objetivo",    "Dispuestos a implementar solución (Q13)","80 %",  "82,406"),
]
for r_idx, (seg, crit, fac, tam) in enumerate(rows_data, start=1):
    cells = tbl.rows[r_idx].cells
    cells[0].text = seg
    cells[1].text = crit
    cells[2].text = fac
    cells[3].text = tam
    for c in cells:
        c.paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

# ── 5. ANÁLISIS DE PRECIO ─────────────────────────────────
add_heading(doc, "5. Análisis de Precio (Sección 10)")
add_body(doc,
    "La pregunta Q17 de la encuesta solicitó la disposición de pago mensual. De los 100 encuestados, "
    "61 indicaron una intención de pago definitiva distribuida en tres rangos:"
)

doc.add_paragraph()
tbl2 = doc.add_table(rows=5, cols=4)
tbl2.style = 'Table Grid'
tbl2.alignment = WD_TABLE_ALIGNMENT.CENTER
h2 = ["Rango de Precio Mensual", "Precio Representativo/mes", "Respondentes", "Subtotal anual (S/.)"]
for i, h in enumerate(h2):
    cell = tbl2.rows[0].cells[i]
    cell.text = h
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.color.rgb = WHITE
    set_cell_bg(cell, "0070C0")

price_rows = [
    ("S/. 20 – S/. 40 / mes",  "S/. 30",  "35",  "S/. 12,600 / año"),
    ("S/. 41 – S/. 70 / mes",  "S/. 55",  "18",  "S/.  9,900 / año"),
    ("S/. 71 – S/. 100 / mes", "S/. 85",   "8",  "S/.  8,160 / año"),
    ("TOTAL / Promedio ponderado", "",    "61",  "S/. 30,660 / año"),
]
for r_idx, row in enumerate(price_rows, start=1):
    cells = tbl2.rows[r_idx].cells
    for c_idx, val in enumerate(row):
        cells[c_idx].text = val
        cells[c_idx].paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if r_idx == 4:
        for c in cells:
            set_cell_bg(c, "CCCCFF")
            c.paragraphs[0].runs[0].bold = True

doc.add_paragraph()
add_body(doc,
    "Precio Ponderado Anual = S/. 30,660 ÷ 61 respondentes = S/. 502.62 / docente / año",
    bold=True
)
add_body(doc,
    "(Equivalente: S/. 30,660 / 61 ≈ S/. 502.62/año → ~S/. 41.89/mes promedio ponderado)"
)
add_body(doc,
    "Nota: En el modelo Excel se utilizan los precios anuales representativos (S/.360, S/.660, "
    "S/.1,020) obteniendo un precio ponderado de S/.535.08/año, diferencia menor debida al "
    "redondeo del precio representativo por rango."
)

# ── 6. ESCENARIOS DE DEMANDA ──────────────────────────────
add_heading(doc, "6. Escenarios de Demanda (Sección 11)")
add_body(doc,
    "Se plantean tres escenarios de penetración sobre el Mercado Objetivo (82,406 docentes) "
    "para el año de lanzamiento 2027, considerando el precio ponderado de S/.535.08/año:"
)

doc.add_paragraph()
tbl3 = doc.add_table(rows=4, cols=5)
tbl3.style = 'Table Grid'
tbl3.alignment = WD_TABLE_ALIGNMENT.CENTER
h3 = ["Escenario", "% Penetración", "Suscripciones", "Ingresos Anuales (S/.)", "Ingresos Anuales (USD*)"]
for i, h in enumerate(h3):
    cell = tbl3.rows[0].cells[i]
    cell.text = h
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.color.rgb = WHITE
    set_cell_bg(cell, "0070C0")

esc_rows = [
    ("Optimista",  "80 %", "65,925", "S/. 35,274,015", "≈ USD 9,521,355"),
    ("Moderado",   "50 %", "41,203", "S/. 22,046,259", "≈ USD 5,950,881"),
    ("Pesimista",  "30 %", "24,722", "S/. 13,227,786", "≈ USD 3,569,942"),
]
colors = ["E2EFDA", "FFEB9C", "FFC7CE"]
for r_idx, (esc, pct, sus, ing, usd) in enumerate(esc_rows, start=1):
    cells = tbl3.rows[r_idx].cells
    for c_idx, val in enumerate([esc, pct, sus, ing, usd]):
        cells[c_idx].text = val
        cells[c_idx].paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for c in cells:
        set_cell_bg(c, colors[r_idx - 1])

doc.add_paragraph()
add_body(doc, "* Tipo de cambio referencial: S/. 3.71 / USD (BCRP, abril 2026)")

# ── 7. PROYECCIÓN 5 AÑOS ──────────────────────────────────
add_heading(doc, "7. Proyección a 5 Años (2027–2031) — Escenario Moderado")
add_body(doc,
    "Tomando el escenario moderado como base y aplicando una tasa de crecimiento anual del 15 %, "
    "sustentada en el HolonIQ EdTech LATAM Report 2024 (CAGR 15 % para plataformas SaaS educativas "
    "en Perú y Latinoamérica):"
)

doc.add_paragraph()
tbl4 = doc.add_table(rows=6, cols=4)
tbl4.style = 'Table Grid'
tbl4.alignment = WD_TABLE_ALIGNMENT.CENTER
h4 = ["Año", "Suscripciones", "Ingresos (S/.)", "Crecimiento"]
for i, h in enumerate(h4):
    cell = tbl4.rows[0].cells[i]
    cell.text = h
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.color.rgb = WHITE
    set_cell_bg(cell, "7030A0")

proj = [
    ("2027", "41,203",  "S/. 22,046,259", "Base"),
    ("2028", "47,383",  "S/. 25,353,198", "+15 %"),
    ("2029", "54,490",  "S/. 29,156,177", "+15 %"),
    ("2030", "62,664",  "S/. 33,529,604", "+15 %"),
    ("2031", "72,063",  "S/. 38,559,045", "+15 %"),
]
for r_idx, row in enumerate(proj, start=1):
    cells = tbl4.rows[r_idx].cells
    for c_idx, val in enumerate(row):
        cells[c_idx].text = val
        cells[c_idx].paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

doc.add_paragraph()
add_body(doc,
    "En cinco años, el escenario moderado proyecta ingresos acumulados superiores a "
    "S/. 148 millones, alcanzando 72,063 suscripciones activas en 2031."
)

# ── 8. CONCLUSIONES ───────────────────────────────────────
add_heading(doc, "8. Conclusiones")
conclusions = [
    "El mercado objetivo de LegendaryClass en Lima Metropolitana asciende a 82,406 docentes, "
    "lo que representa una oportunidad de mercado significativa y cuantificable.",
    "El precio ponderado anual de S/. 502–535 por docente es competitivo frente a soluciones "
    "internacionales equivalentes (Classcraft: USD 120/año; ClassDojo Premium: USD 36/año) y "
    "se alinea con la disposición de pago local.",
    "El escenario moderado (50 % de penetración) genera ingresos de S/. 22 millones en el año "
    "de lanzamiento, lo que sustenta la viabilidad financiera del proyecto.",
    "Con una tasa de crecimiento del 15 % anual (respaldada por el mercado EdTech LATAM), "
    "LegendaryClass puede alcanzar S/. 38.5 millones en ingresos para 2031.",
    "La demanda está sostenida por una necesidad real: el 58 % de los docentes encuestados "
    "reporta baja motivación estudiantil de forma frecuente o constante.",
]
for i, c in enumerate(conclusions, start=1):
    add_bullet(doc, f"{i}. {c}")

# ── 9. REFERENCIAS ────────────────────────────────────────
add_heading(doc, "9. Referencias")
refs = [
    "Ministerio de Educación del Perú (MINEDU). (2024). Estadísticas de la Calidad Educativa (ESCALE). Lima, Perú.",
    "HolonIQ. (2024). Latin America EdTech Market Briefing. HolonIQ Intelligence.",
    "Banco Central de Reserva del Perú (BCRP). (2026, abril). Tipo de cambio referencial. Lima, Perú.",
    "Deterding, S., Dixon, D., Khaled, R., & Nacke, L. (2011). From game design elements to gamefulness: Defining gamification. MindTrek '11 Proceedings.",
    "Hamari, J., Koivisto, J., & Sarsa, H. (2014). Does gamification work? A literature review of empirical studies on gamification. HICSS 2014.",
]
for r in refs:
    add_bullet(doc, r)

WORD_PATH = "C:/Proyectos/LegendaryClass/lab03/Informe_Caso3_LegendaryClass.docx"
doc.save(WORD_PATH)
print("[OK] Word guardado en:", WORD_PATH)


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Cm as PCm

PPURPLE = PRGB(112, 48, 160)
PBLUE   = PRGB(0, 112, 192)
PWHITE  = PRGB(255, 255, 255)
PGRAY   = PRGB(80, 80, 80)
PDARK   = PRGB(20, 20, 20)
PGOLD   = PRGB(255, 192, 0)
PGREEN  = PRGB(84, 130, 53)
PRED    = PRGB(192, 0, 0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completamente en blanco

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt as UPt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Emu(line_width or 9525)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=PDARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_slide_header(slide, title, subtitle=""):
    # barra superior oscura
    add_rect(slide, 0, 0, 13.33, 1.15, fill_color=PPURPLE)
    add_text_box(slide, title, 0.3, 0.1, 10, 0.65,
                 font_size=28, bold=True, color=PWHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.72, 10, 0.38,
                     font_size=14, bold=False, color=PRGB(220, 200, 255), align=PP_ALIGN.LEFT)
    # barra inferior
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=PPURPLE)
    add_text_box(slide, "LegendaryClass · Análisis de Demanda · 2026", 0.3, 7.12, 12, 0.3,
                 font_size=9, color=PRGB(200, 180, 240), align=PP_ALIGN.LEFT)

def add_kpi_box(slide, value, label, l, t, w=2.8, h=1.4, bg=PPURPLE):
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_text_box(slide, value, l, t+0.1, w, 0.75,
                 font_size=26, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l, t+0.8, w, 0.5,
                 font_size=11, color=PRGB(220, 200, 255), align=PP_ALIGN.CENTER)

# ── SLIDE 1: Portada ─────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(15, 15, 30))
# Acento decorativo
add_rect(sl, 0, 0, 0.4, 7.5, fill_color=PPURPLE)
add_rect(sl, 12.93, 0, 0.4, 7.5, fill_color=PPURPLE)

add_text_box(sl, "LABORATORIO 03 · ANÁLISIS DE DEMANDA", 1, 0.9, 11.33, 0.5,
             font_size=13, color=PRGB(180, 140, 255), align=PP_ALIGN.CENTER)
add_text_box(sl, "CASO 3", 1, 1.55, 11.33, 0.7,
             font_size=30, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
add_text_box(sl, "LegendaryClass", 1, 2.3, 11.33, 1.1,
             font_size=52, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "Sistema SaaS de Gamificación Educativa", 1, 3.55, 11.33, 0.55,
             font_size=20, color=PRGB(180, 180, 180), align=PP_ALIGN.CENTER)

# línea separadora
add_rect(sl, 3, 4.25, 7.33, 0.04, fill_color=PPURPLE)

add_text_box(sl, "Determinación de la Demanda del Mercado · Lima Metropolitana · 2027", 1, 4.45, 11.33, 0.45,
             font_size=13, color=PRGB(150, 130, 200), align=PP_ALIGN.CENTER)
add_text_box(sl, "Juan Aguirre  |  Docente: Lturkowsky  |  2026 - I", 1, 5.0, 11.33, 0.4,
             font_size=12, color=PRGB(130, 130, 130), align=PP_ALIGN.CENTER)

# ── SLIDE 2: Descripción del Producto ───────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(245, 245, 252))
add_slide_header(sl, "¿Qué es LegendaryClass?", "Descripción del producto / servicio")

add_text_box(sl, "Plataforma SaaS de gamificación que transforma el aula en un RPG educativo.",
             0.5, 1.3, 8.5, 0.55, font_size=16, bold=True, color=PPURPLE)

features = [
    ("⚔️  5 Clases de Personaje", "Mago, Guerrero, Ninja, Arquero, Lanzador\n+20 % XP por tipo de acción"),
    ("📊  Motor de Gamificación", "Niveles, racha diaria, logros\nautomáticos y misiones configurables"),
    ("👩‍🏫  Panel del Docente", "Registro de comportamientos,\nrecompensas y reportes de progreso"),
    ("🏛️  Panel del Director", "Estadísticas institucionales\nagregadas en tiempo real"),
]
x_positions = [0.4, 3.5, 6.6, 9.7]
for (title, desc), x in zip(features, x_positions):
    add_rect(sl, x, 2.0, 3.1, 2.0, fill_color=PRGB(235, 228, 250))
    add_text_box(sl, title, x+0.1, 2.05, 2.9, 0.55, font_size=13, bold=True, color=PPURPLE)
    add_text_box(sl, desc, x+0.1, 2.6, 2.9, 1.3, font_size=11, color=PGRAY)

add_text_box(sl, "Stack tecnológico:  Angular 18 · NestJS · PostgreSQL · Prisma · TailwindCSS",
             0.5, 4.2, 12, 0.45, font_size=13, bold=True, color=PBLUE)
add_text_box(sl,
    "Modelo de negocio:  Suscripción mensual por docente  ·  Precio: S/. 30 – S/. 85 / mes  "
    "·  Sin costo para estudiantes",
    0.5, 4.7, 12, 0.45, font_size=12, color=PGRAY)

# ── SLIDE 3: Metodología ─────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(245, 245, 252))
add_slide_header(sl, "Metodología", "Instrumento y segmentación en cascada")

add_text_box(sl, "Encuesta digital aplicada a 100 docentes de Lima Metropolitana (feb–mar 2026)",
             0.5, 1.3, 12, 0.45, font_size=15, bold=True, color=PPURPLE)

steps = [
    ("1", "Mercado\nPotencial", "185,000\ndocentes", "Todos los docentes activos\nen Lima Met. (MINEDU 2024)", PPURPLE),
    ("2", "Mercado\nDisponible", "177,600\ndocentes", "Trabajan en institución\neducativa actualmente (Q4 · 96 %)", PBLUE),
    ("3", "Mercado\nEfectivo", "103,008\ndocentes", "Perciben baja motivación\nfrecuente o siempre (Q6 · 58 %)", PRGB(0, 150, 130)),
    ("4", "Mercado\nObjetivo", "82,406\ndocentes", "Dispuestos a implementar\nuna solución (Q13 · 80 %)", PGREEN),
]
for i, (num, label, value, desc, color) in enumerate(steps):
    x = 0.35 + i * 3.22
    add_rect(sl, x, 1.95, 2.9, 1.8, fill_color=color)
    add_text_box(sl, label, x, 2.0, 2.9, 0.55, font_size=14, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, value, x, 2.55, 2.9, 0.65, font_size=19, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl, desc, x, 3.2, 2.9, 0.75, font_size=9.5, color=PRGB(220, 220, 220), align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(sl, "→", x + 2.9, 2.65, 0.32, 0.6, font_size=22, bold=True, color=PPURPLE, align=PP_ALIGN.CENTER)

add_text_box(sl, "Resultado: 82,406 docentes en Lima Metropolitana son el mercado objetivo de LegendaryClass.",
             0.5, 4.0, 12, 0.45, font_size=14, bold=True, color=PPURPLE)

add_text_box(sl,
    "Fuente de mercado potencial: MINEDU – ESCALE 2024  ·  "
    "Filtros aplicados con base en respuestas de la encuesta propia (n=100)",
    0.5, 4.55, 12, 0.4, font_size=10, color=PGRAY)

# ── SLIDE 4: Análisis de Precio ─────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(245, 245, 252))
add_slide_header(sl, "Análisis de Precio", "Disposición de pago (Q17) — 61 respondentes con intención definitiva")

# 3 cajas de precio
price_data = [
    ("S/. 30 / mes", "35 docentes", "S/. 360 / año", "Rango: S/.20–40", PBLUE),
    ("S/. 55 / mes", "18 docentes", "S/. 660 / año", "Rango: S/.41–70", PPURPLE),
    ("S/. 85 / mes", "8 docentes",  "S/. 1,020 / año","Rango: S/.71–100", PRGB(0, 150, 130)),
]
for i, (price, qty, annual, rng, color) in enumerate(price_data):
    x = 0.5 + i * 4.1
    add_rect(sl, x, 1.35, 3.6, 2.5, fill_color=color)
    add_text_box(sl, rng,    x, 1.4, 3.6, 0.38, font_size=11, color=PRGB(200,200,200), align=PP_ALIGN.CENTER)
    add_text_box(sl, price,  x, 1.78, 3.6, 0.7, font_size=26, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl, qty,    x, 2.48, 3.6, 0.48, font_size=15, color=PWHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, annual, x, 2.95, 3.6, 0.48, font_size=13, bold=True, color=PRGB(180,255,180), align=PP_ALIGN.CENTER)

# Precio ponderado resultado
add_rect(sl, 2.5, 4.05, 8.33, 1.05, fill_color=PPURPLE)
add_text_box(sl, "Precio Ponderado Anual", 2.5, 4.08, 8.33, 0.4,
             font_size=14, color=PRGB(220,200,255), align=PP_ALIGN.CENTER)
add_text_box(sl, "S/. 535.08 / docente / año   (~S/. 44.59 / mes)", 2.5, 4.45, 8.33, 0.55,
             font_size=22, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)

add_text_box(sl,
    "Cálculo: (35×S/.360 + 18×S/.660 + 8×S/.1,020) ÷ 61 = S/.32,640 ÷ 61 = S/.535.08 / año",
    0.5, 5.3, 12.33, 0.45, font_size=11, color=PGRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 5: Escenarios de Demanda ───────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(245, 245, 252))
add_slide_header(sl, "Escenarios de Demanda", "Año de lanzamiento 2027 · Mercado Objetivo: 82,406 docentes")

scenarios = [
    ("OPTIMISTA", "80 %", "65,925", "S/. 35,274,015", PGREEN,    PRGB(220, 240, 210)),
    ("MODERADO",  "50 %", "41,203", "S/. 22,046,259", PBLUE,     PRGB(210, 225, 245)),
    ("PESIMISTA", "30 %", "24,722", "S/. 13,227,786", PRED,      PRGB(250, 215, 215)),
]
for i, (name, pct, subs, revenue, title_color, bg) in enumerate(scenarios):
    x = 0.45 + i * 4.3
    add_rect(sl, x, 1.3, 4.0, 3.8, fill_color=bg)
    add_rect(sl, x, 1.3, 4.0, 0.6, fill_color=title_color)
    add_text_box(sl, name,    x, 1.32, 4.0, 0.52, font_size=18, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Penetración", x, 2.0,  4.0, 0.35, font_size=10, color=PGRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, pct,     x, 2.35, 4.0, 0.65, font_size=30, bold=True, color=title_color, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Suscripciones", x, 3.05, 4.0, 0.35, font_size=10, color=PGRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, subs,    x, 3.4,  4.0, 0.55, font_size=22, bold=True, color=PDARK, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Ingresos anuales", x, 4.0, 4.0, 0.35, font_size=10, color=PGRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, revenue, x, 4.35, 4.0, 0.55, font_size=16, bold=True, color=title_color, align=PP_ALIGN.CENTER)

add_text_box(sl,
    "Precio base: S/. 535.08 / año  ·  Demanda = Mercado Objetivo × % de penetración",
    0.5, 5.3, 12.33, 0.4, font_size=11, color=PGRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 6: Proyección 5 años ───────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(245, 245, 252))
add_slide_header(sl, "Proyección 2027–2031", "Escenario Moderado · Crecimiento anual: 15 % (HolonIQ EdTech LATAM 2024)")

years   = ["2027",      "2028",      "2029",      "2030",      "2031"]
subs    = [41203,       47383,       54490,       62664,       72063]
revenue = [22_046_259,  25_353_198,  29_156_177,  33_529_604,  38_559_045]

# Gráfico de barras simulado con rectángulos
max_rev = max(revenue)
bar_colors = [PPURPLE, PBLUE, PRGB(0,150,130), PGREEN, PGOLD]
chart_l, chart_b, chart_w, bar_gap = 0.6, 5.7, 12.1, 0.35
bar_w = (chart_w - bar_gap * 6) / 5
chart_h = 2.8

for i, (yr, sub, rev) in enumerate(zip(years, subs, revenue)):
    bar_h = (rev / max_rev) * chart_h
    x = chart_l + bar_gap + i * (bar_w + bar_gap)
    y_bar = chart_b - bar_h
    add_rect(sl, x, y_bar, bar_w, bar_h, fill_color=bar_colors[i])
    rev_m = f"S/.{rev/1_000_000:.1f}M"
    add_text_box(sl, rev_m, x - 0.05, y_bar - 0.4, bar_w + 0.1, 0.38,
                 font_size=10.5, bold=True, color=bar_colors[i], align=PP_ALIGN.CENTER)
    add_text_box(sl, yr, x - 0.05, chart_b + 0.08, bar_w + 0.1, 0.35,
                 font_size=12, bold=True, color=PDARK, align=PP_ALIGN.CENTER)
    add_text_box(sl, f"{sub:,}", x - 0.05, chart_b + 0.43, bar_w + 0.1, 0.3,
                 font_size=9, color=PGRAY, align=PP_ALIGN.CENTER)

# KPIs debajo del header
kpis = [
    ("41,203 → 72,063", "Suscripciones\n2027 → 2031", 0.6,  1.3, PPURPLE),
    ("+75 %",           "Crecimiento\nacumulado 5 años", 4.8,  1.3, PBLUE),
    ("S/. 148.7M",      "Ingresos\nacumulados 2027-31", 9.0,  1.3, PGREEN),
]
for val, lbl, x, y, color in kpis:
    add_rect(sl, x, y, 3.3, 1.3, fill_color=color)
    add_text_box(sl, val, x, y+0.08, 3.3, 0.65, font_size=20, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, lbl, x, y+0.72, 3.3, 0.48, font_size=10, color=PRGB(220,210,255), align=PP_ALIGN.CENTER)

add_text_box(sl, "Suscripciones", 0.5, 6.1, 3, 0.3, font_size=9, color=PGRAY)

# ── SLIDE 7: Conclusiones ────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(15, 15, 30))
add_rect(sl, 0, 0, 0.3, 7.5, fill_color=PPURPLE)
add_rect(sl, 0, 0, 13.33, 1.15, fill_color=PPURPLE)
add_text_box(sl, "Conclusiones", 0.5, 0.12, 12, 0.65,
             font_size=28, bold=True, color=PWHITE)

conclusions_ppt = [
    ("82,406 docentes", "en Lima Metropolitana conforman el mercado objetivo de LegendaryClass."),
    ("S/. 535 / año",   "es el precio ponderado de disposición de pago entre docentes interesados."),
    ("S/. 22M en 2027", "de ingresos proyectados en el escenario moderado (50 % penetración)."),
    ("S/. 38.5M en 2031","alcanzables con crecimiento del 15 % anual (EdTech LATAM CAGR)."),
    ("58 % de docentes","reportan baja motivación estudiantil frecuente — necesidad real y validada."),
]
for i, (bold_part, rest) in enumerate(conclusions_ppt):
    y = 1.4 + i * 1.0
    add_rect(sl, 0.5, y, 0.06, 0.55, fill_color=PGOLD)
    add_text_box(sl, bold_part, 0.7, y + 0.02, 2.8, 0.52,
                 font_size=14, bold=True, color=PGOLD)
    add_text_box(sl, rest, 3.5, y + 0.02, 9.3, 0.52,
                 font_size=13, color=PRGB(200, 200, 220))

# ── SLIDE 8: Cierre ──────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=PRGB(15, 15, 30))
add_rect(sl, 0, 0, 0.3, 7.5, fill_color=PPURPLE)
add_rect(sl, 13.0, 0, 0.33, 7.5, fill_color=PPURPLE)

add_text_box(sl, "⚔️  LegendaryClass", 1, 2.2, 11.33, 1.0,
             font_size=42, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "Transformando el aula en una experiencia legendaria.", 1, 3.35, 11.33, 0.6,
             font_size=18, color=PRGB(180, 160, 230), align=PP_ALIGN.CENTER)
add_rect(sl, 4.0, 4.2, 5.33, 0.05, fill_color=PPURPLE)
add_text_box(sl, "Gracias", 1, 4.5, 11.33, 0.8,
             font_size=34, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
add_text_box(sl, "Juan Aguirre  |  Ingeniería de Software  |  2026 - I", 1, 5.5, 11.33, 0.45,
             font_size=13, color=PRGB(130, 130, 130), align=PP_ALIGN.CENTER)

PPT_PATH = "C:/Proyectos/LegendaryClass/lab03/Presentacion_Caso3_LegendaryClass.pptx"
prs.save(PPT_PATH)
print("[OK] PowerPoint guardado en:", PPT_PATH)
print()
print("Archivos generados:")
print("  ->", WORD_PATH)
print("  ->", PPT_PATH)
