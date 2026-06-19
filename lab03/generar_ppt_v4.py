"""
PPT Caso 3 - LegendaryClass
Estilo fiel a la web PHP original:
  - Fondo: imagen fondo.png con overlay blanco
  - Título: dorado #fbbf24 / #f59e0b, serif (Cinzel-style)
  - Header barra: dark navy→purple con borde dorado inferior
  - Acento: dorado/ámbar, no morado saturado
"""
import sys, io, shutil, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN

# ── Paleta extraída directamente del CSS original ────────
NAVY       = RGB( 15,  23,  42)   # #0f172a — fondo header
PURPLE_D   = RGB( 88,  28, 135)   # #581c87 — purple-900
PURPLE_M   = RGB(124,  58, 237)   # #7c3aed — purple-600
GOLD_400   = RGB(251, 191,  36)   # #fbbf24 — título principal
GOLD_500   = RGB(245, 158,  11)   # #f59e0b — acento
GOLD_600   = RGB(217, 119,   6)   # #d97706 — sombra dorada
GOLD_800   = RGB(146,  64,  14)   # #92400e — dorado oscuro
WHITE      = RGB(255, 255, 255)
CREAM      = RGB(253, 251, 244)   # fondo slides contenido
GRAY_100   = RGB(243, 244, 246)
GRAY_300   = RGB(209, 213, 219)
GRAY_500   = RGB(107, 114, 128)
GRAY_700   = RGB( 55,  65,  81)
GRAY_800   = RGB( 31,  41,  55)   # texto oscuro principal
PURPLE_L   = RGB(196, 181, 253)   # purple-300

FONDO = "C:/xampp/htdocs/LegendaryClass/public/fondo.png"

# ── Datos ────────────────────────────────────────────────
INSTITUCION = "TECSUP"
CARRERA     = "Ingeniería de Software  ·  V Ciclo"
CURSO       = "Diseño de Proyectos de Innovación"
DOCENTE     = "Turkowsky Vizcarra, Luisa"
ALUMNO      = "Aguirre, Juan"
ANIO        = "2026 - I"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────
def _rect(sl, l, t, w, h, fill=None, line_color=None, line_w=9525):
    sp = sl.shapes.add_shape(1,
        Inches(max(l, -1)), Inches(max(t, -1)),
        Inches(min(w, 15)),  Inches(min(h, 10)))
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color; sp.line.width = Emu(line_w)
    else:
        sp.line.fill.background()
    return sp

def txt(sl, text, l, t, w, h,
        size=13, bold=False, color=GRAY_800,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def bg_content(sl):
    """Fondo claro cremoso para slides de contenido"""
    _rect(sl, 0, 0, 13.33, 7.5, fill=CREAM)
    # Toque sutil de color en esquinas
    _rect(sl, 0, 0, 2.5, 2.5, fill=RGB(255,253,235))
    _rect(sl, 10.83, 5, 2.5, 2.5, fill=RGB(255,253,235))

def epic_header(sl, title, subtitle=""):
    """Barra header idéntica al epic-header del PHP"""
    # Fondo oscuro navy-purple
    _rect(sl, 0, 0, 13.33, 1.35, fill=NAVY)
    # Degradado simulado con capa encima
    _rect(sl, 2.5, 0, 8.33, 1.35, fill=PURPLE_D)
    _rect(sl, 4.5, 0, 4.33, 1.35, fill=PURPLE_M)
    # Borde inferior dorado (como border-image del CSS)
    _rect(sl, 0, 1.3, 3.33, 0.07, fill=GOLD_400)
    _rect(sl, 3.33, 1.3, 3.33, 0.07, fill=GOLD_500)
    _rect(sl, 6.66, 1.3, 3.33, 0.07, fill=GOLD_600)
    _rect(sl, 10.0, 1.3, 3.33, 0.07, fill=GOLD_800)
    # Título en dorado — estilo Cinzel
    txt(sl, title.upper(), 0.4, 0.1, 12.5, 0.72,
        size=22, bold=True, color=GOLD_400, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(sl, subtitle, 0.4, 0.8, 12.5, 0.42,
            size=11, color=PURPLE_L, italic=False)
    # Footer
    _rect(sl, 0, 7.3, 13.33, 0.2, fill=NAVY)
    _rect(sl, 0, 7.3, 13.33, 0.04, fill=GOLD_500)
    txt(sl, f"LegendaryClass  ·  {CURSO}  ·  {INSTITUCION}  ·  {ANIO}",
        0.4, 7.35, 12.5, 0.2, size=8, color=PURPLE_L)

def card(sl, l, t, w, h, border_color=GOLD_500):
    """Card blanca con borde dorado"""
    _rect(sl, l, t, w, h, fill=WHITE, line_color=border_color, line_w=12700)

def gold_line(sl, l, t, w):
    _rect(sl, l, t, w, 0.05, fill=GOLD_500)

def section_tag(sl, text, l, t):
    """Etiqueta de sección dorada"""
    w = len(text) * 0.092 + 0.3
    _rect(sl, l, t, w, 0.32, fill=NAVY)
    txt(sl, text.upper(), l+0.08, t+0.02, w-0.1, 0.28,
        size=9, bold=True, color=GOLD_400, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Fondo: imagen fondo.png si existe
if os.path.exists(FONDO):
    from pptx.util import Inches as _I
    pic = sl.shapes.add_picture(FONDO, _I(0), _I(0), _I(13.33), _I(7.5))
    pic.zorder = 0
# Overlay blanco semitransparente (simula el CSS rgba(255,255,255,0.82))
overlay = _rect(sl, 0, 0, 13.33, 7.5, fill=RGB(255, 255, 255))
from pptx.oxml.ns import qn
from lxml import etree
# Hacer overlay semitransparente via XML
sp_elem = overlay._element
spPr = sp_elem.find(qn('p:spPr'))
solidFill = spPr.find('.//' + qn('a:solidFill'))
if solidFill is not None:
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', 'FFFFFF')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '82000')  # 82% opacity

# Header épico
_rect(sl, 0, 0, 13.33, 1.6, fill=NAVY)
_rect(sl, 2.0, 0, 9.33, 1.6, fill=PURPLE_D)
_rect(sl, 4.5, 0, 4.33, 1.6, fill=PURPLE_M)
_rect(sl, 0, 1.55, 3.33, 0.08, fill=GOLD_400)
_rect(sl, 3.33, 1.55, 3.33, 0.08, fill=GOLD_500)
_rect(sl, 6.66, 1.55, 3.33, 0.08, fill=GOLD_600)
_rect(sl, 10.0, 1.55, 3.33, 0.08, fill=GOLD_800)

# Título principal — LEGENDARYCLASS en dorado
txt(sl, "LEGENDARYCLASS", 0.5, 0.18, 12.33, 1.0,
    size=52, bold=True, color=GOLD_400, align=PP_ALIGN.CENTER)

# Panel central
card(sl, 2.8, 2.0, 7.73, 4.6)
txt(sl, "⚔️  Bienvenido al Portal de Aventureros  ⚔️",
    3.0, 2.2, 7.33, 0.5, size=13, color=GRAY_500, align=PP_ALIGN.CENTER, italic=True)
gold_line(sl, 3.2, 2.82, 6.93)
txt(sl, CURSO, 3.0, 2.95, 7.33, 0.45, size=13, bold=True, color=GRAY_800, align=PP_ALIGN.CENTER)
txt(sl, "Caso 3 — Parte Práctica", 3.0, 3.45, 7.33, 0.45,
    size=15, bold=True, color=GOLD_600, align=PP_ALIGN.CENTER)
gold_line(sl, 3.2, 4.0, 6.93)
txt(sl, f"Docente:  {DOCENTE}", 3.0, 4.1, 7.33, 0.38,
    size=11, color=GRAY_700, align=PP_ALIGN.CENTER)
txt(sl, f"Alumno:   {ALUMNO}", 3.0, 4.5, 7.33, 0.38,
    size=11, color=GRAY_700, align=PP_ALIGN.CENTER)
txt(sl, f"{INSTITUCION}  ·  {CARRERA}  ·  {ANIO}",
    3.0, 4.98, 7.33, 0.38, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)
gold_line(sl, 3.2, 5.45, 6.93)
txt(sl, "Determinación de la Demanda del Mercado  ·  Lima Metropolitana  ·  2027",
    3.0, 5.55, 7.33, 0.38, size=10, color=GOLD_600, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════
# SLIDE 2 — DESCRIPCIÓN DEL PRODUCTO
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Descripción del Producto", "LegendaryClass — Sistema SaaS de Gamificación Educativa")

txt(sl, "Transforma el aula en un RPG educativo. Los estudiantes eligen clase de personaje, "
        "acumulan XP, suben de nivel y canjean recompensas. El docente gestiona comportamientos, misiones y reportes.",
    0.45, 1.52, 12.43, 0.72, size=12.5, color=GRAY_700, wrap=True)

features = [
    ("⚔️", "5 Clases RPG",      "Mago · Guerrero · Ninja\nArquero · Lanzador\n+20 % XP por tipo de acción"),
    ("📈", "Motor de XP",        "Niveles progresivos, racha\ndiaria y logros automáticos"),
    ("👩‍🏫","Panel Docente",     "Comportamientos, recompensas\ny reportes por aula"),
    ("🏛️","Panel Director",     "Estadísticas globales\nde la institución"),
]
for i, (ico, tit, desc) in enumerate(features):
    x = 0.45 + i * 3.23
    card(sl, x, 2.42, 3.05, 2.85)
    _rect(sl, x, 2.42, 3.05, 0.06, fill=GOLD_500)
    txt(sl, ico,  x+0.15, 2.52, 0.65, 0.55, size=26)
    txt(sl, tit,  x+0.15, 3.1,  2.75, 0.45, size=12.5, bold=True, color=GOLD_600)
    txt(sl, desc, x+0.15, 3.56, 2.75, 1.55, size=10.5, color=GRAY_500, wrap=True)

_rect(sl, 0.45, 5.45, 12.43, 0.05, fill=GOLD_500)
txt(sl, "Stack:  Angular 18  ·  NestJS  ·  PostgreSQL  ·  Prisma  ·  TailwindCSS",
    0.45, 5.57, 8.5, 0.38, size=11, bold=True, color=NAVY)
txt(sl, "Precio:  S/. 30 – 85 / mes por docente   ·   Sin costo para estudiantes",
    0.45, 5.97, 10, 0.35, size=11, color=GRAY_500)


# ════════════════════════════════════════════════════════
# SLIDE 3 — METODOLOGÍA Y SEGMENTACIÓN
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Metodología", "Encuesta a 100 docentes · Lima Metropolitana · feb–mar 2026")

# Encuesta info
card(sl, 0.45, 1.52, 5.7, 1.65)
_rect(sl, 0.45, 1.52, 5.7, 0.06, fill=GOLD_500)
txt(sl, "Encuesta de campo", 0.65, 1.62, 5.3, 0.42, size=12.5, bold=True, color=GOLD_600)
for j, item in enumerate([
    "100 docentes · Lima Metropolitana",
    "Google Forms · 20 preguntas",
    "Periodo: febrero – marzo 2026",
]):
    txt(sl, "·  " + item, 0.65, 2.1+j*0.36, 5.3, 0.32, size=11.5, color=GRAY_700)

# Preguntas clave
card(sl, 6.45, 1.52, 6.43, 1.65)
_rect(sl, 6.45, 1.52, 6.43, 0.06, fill=GOLD_500)
txt(sl, "Preguntas clave", 6.65, 1.62, 6.1, 0.42, size=12.5, bold=True, color=GOLD_600)
qs = ["Q4  — Situación laboral  →  Disponible",
      "Q6  — Baja motivación  →  Efectivo",
      "Q13 — Disposición gamificación  →  Objetivo",
      "Q17 — Disposición de pago  →  Precio"]
for j, q in enumerate(qs):
    txt(sl, q, 6.65, 2.1+j*0.36, 6.1, 0.32, size=11, color=GRAY_700)

# Cascada
txt(sl, "Segmentación en Cascada", 0.45, 3.35, 8, 0.42, size=13, bold=True, color=NAVY)
cascade = [
    ("POTENCIAL",   "185,000",  "MINEDU 2024",   NAVY,     GOLD_400),
    ("DISPONIBLE",  "177,600",  "× 96 %  (Q4)",  PURPLE_D, GOLD_400),
    ("EFECTIVO",    "103,008",  "× 58 %  (Q6)",  PURPLE_M, GOLD_400),
    ("OBJETIVO",     "82,406",  "× 80 %  (Q13)", GOLD_600, WHITE),
]
for i, (lbl, val, crit, bg_col, fg) in enumerate(cascade):
    x = 0.45 + i * 3.23
    _rect(sl, x, 3.88, 3.05, 2.3, fill=bg_col)
    _rect(sl, x, 3.88, 3.05, 0.06, fill=GOLD_400)
    txt(sl, lbl,  x+0.15, 3.97, 2.75, 0.38, size=10, bold=True, color=fg)
    txt(sl, val,  x+0.15, 4.38, 2.75, 0.65, size=26, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(sl, crit, x+0.15, 5.05, 2.75, 0.35, size=10.5, color=RGB(200,180,150) if bg_col!=GOLD_600 else NAVY,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "›", x+3.05, 4.55, 0.18, 0.5, size=20, bold=True, color=GOLD_500, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS DE ENCUESTA
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Resultados de la Encuesta", "Preguntas clave para la segmentación  ·  n = 100")

# Barras Q6
txt(sl, "Q6  —  Frecuencia de baja motivación estudiantil",
    0.45, 1.52, 7.5, 0.42, size=13, bold=True, color=NAVY)

q6 = [("Siempre",        22, GOLD_600,  True),
      ("Frecuentemente",  36, GOLD_500,  True),
      ("A veces",         30, GRAY_300,  False),
      ("Raramente",        8, GRAY_300,  False),
      ("Nunca",            4, GRAY_300,  False)]
for j, (lbl, val, col, highlight) in enumerate(q6):
    y = 2.05 + j * 0.72
    _rect(sl, 0.45, y, 7.0, 0.58, fill=GRAY_100)
    bw = (val/36)*6.4
    _rect(sl, 0.45, y, bw, 0.58, fill=col)
    txt(sl, lbl, 0.58, y+0.13, 2.2, 0.3,
        size=11, bold=highlight, color=WHITE if highlight else GRAY_700)
    txt(sl, f"{val} %", 0.45+bw+0.1, y+0.13, 0.65, 0.3,
        size=12, bold=highlight, color=col if col!=GRAY_300 else GRAY_500)
    if highlight:
        txt(sl, "← incluido en filtro", 0.45+bw+0.8, y+0.15, 2.0, 0.28,
            size=9, color=GOLD_600, italic=True)

# Resultado 58%
_rect(sl, 0.45, 5.65, 7.0, 0.72, fill=NAVY)
txt(sl, "Siempre + Frecuentemente  =  58 %   →   Mercado Efectivo",
    0.6, 5.78, 6.6, 0.42, size=12, bold=True, color=GOLD_400)

# Panel derecho
txt(sl, "Resumen por pregunta", 7.8, 1.52, 5.1, 0.42, size=13, bold=True, color=NAVY)
hallazgos = [
    ("Q4",  "96 %",  "Trabajan en institución educativa",    NAVY),
    ("Q6",  "58 %",  "Baja motivación frecuente/siempre",    PURPLE_D),
    ("Q13", "80 %",  "Implementarían gamificación",          GOLD_600),
    ("Q17", "61/100","Con intención definitiva de pago",     GRAY_700),
]
for j, (q, val, desc, col) in enumerate(hallazgos):
    y = 2.05 + j * 1.12
    card(sl, 7.8, y, 5.1, 0.95)
    _rect(sl, 7.8, y, 5.1, 0.06, fill=GOLD_500)
    txt(sl, q,   7.95, y+0.12, 0.6, 0.38, size=10, bold=True, color=GOLD_600)
    txt(sl, val, 8.6,  y+0.08, 1.8, 0.5,  size=22, bold=True, color=col)
    txt(sl, desc, 10.4, y+0.18, 2.3, 0.5, size=10, color=GRAY_500, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 5 — ANÁLISIS DE PRECIO
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Análisis de Precio", "Q17 — Disposición de pago  ·  61 respondentes con intención definitiva")

price_cards = [
    ("S/. 20 – 40 / mes",  "S/. 30",  "S/. 360 / año",  35, NAVY),
    ("S/. 41 – 70 / mes",  "S/. 55",  "S/. 660 / año",  18, PURPLE_D),
    ("S/. 71 – 100 / mes", "S/. 85",  "S/. 1,020 / año", 8, PURPLE_M),
]
for i, (rng, mo, yr, qty, col) in enumerate(price_cards):
    x = 0.45 + i * 4.15
    card(sl, x, 1.52, 3.9, 3.6)
    _rect(sl, x, 1.52, 3.9, 0.06, fill=GOLD_500)
    txt(sl, rng, x+0.15, 1.62, 3.6, 0.38, size=10.5, color=GRAY_500, align=PP_ALIGN.CENTER)
    txt(sl, mo,  x+0.15, 2.05, 3.6, 0.82, size=38, bold=True, color=col, align=PP_ALIGN.CENTER)
    gold_line(sl, x+0.3, 2.95, 3.3)
    txt(sl, yr,  x+0.15, 3.05, 3.6, 0.5, size=18, bold=True, color=GOLD_600, align=PP_ALIGN.CENTER)
    txt(sl, f"{qty} docentes", x+0.15, 3.62, 3.6, 0.38,
        size=12, color=GRAY_500, align=PP_ALIGN.CENTER)

# Resultado precio ponderado
_rect(sl, 0.45, 5.3, 12.43, 1.5, fill=NAVY)
_rect(sl, 0.45, 5.3, 12.43, 0.07, fill=GOLD_400)
_rect(sl, 0.45, 6.73, 12.43, 0.07, fill=GOLD_400)
txt(sl, "Precio Ponderado Anual", 0.7, 5.38, 6, 0.42, size=11.5, color=GOLD_400)
txt(sl, "S/. 535.08 / docente / año   ≈   S/. 44.59 / mes",
    0.7, 5.8, 10, 0.65, size=28, bold=True, color=WHITE)
txt(sl, "( 35 × S/.360  +  18 × S/.660  +  8 × S/.1,020 )  ÷  61  =  S/.32,640 ÷ 61",
    0.7, 6.46, 11.5, 0.32, size=10, color=PURPLE_L, italic=True)


# ════════════════════════════════════════════════════════
# SLIDE 6 — ESCENARIOS
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Escenarios de Demanda", "Año 2027  ·  Mercado Objetivo: 82,406 docentes  ·  Precio: S/. 535.08 / año")

esc = [
    ("Optimista",  "80 %",  "65,925",  "S/. 35,274,015", NAVY,     GOLD_400),
    ("Moderado",   "50 %",  "41,203",  "S/. 22,046,259", PURPLE_D, GOLD_400),
    ("Pesimista",  "30 %",  "24,722",  "S/. 13,227,786", GRAY_700, GOLD_400),
]
for i, (name, pct, subs, rev, bg_col, fg) in enumerate(esc):
    x = 0.45 + i * 4.3
    _rect(sl, x, 1.52, 4.08, 5.42, fill=bg_col)
    _rect(sl, x, 1.52, 4.08, 0.07, fill=GOLD_400)
    _rect(sl, x, 6.87, 4.08, 0.07, fill=GOLD_400)
    txt(sl, name,  x+0.2, 1.62, 3.7, 0.48, size=16, bold=True, color=GOLD_400)
    txt(sl, pct,   x+0.2, 2.15, 3.7, 1.1,  size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "penetración", x+0.2, 3.28, 3.7, 0.35, size=10, color=PURPLE_L, align=PP_ALIGN.CENTER)
    _rect(sl, x+0.3, 3.72, 3.48, 0.04, fill=GOLD_600)
    txt(sl, "Suscripciones", x+0.2, 3.84, 3.7, 0.35, size=10, color=PURPLE_L, align=PP_ALIGN.CENTER)
    txt(sl, subs,  x+0.2, 4.2,  3.7, 0.62, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(sl, x+0.3, 4.9, 3.48, 0.04, fill=GOLD_600)
    txt(sl, "Ingresos anuales", x+0.2, 5.02, 3.7, 0.35, size=10, color=PURPLE_L, align=PP_ALIGN.CENTER)
    txt(sl, rev,   x+0.2, 5.4,  3.7, 0.52, size=15, bold=True, color=GOLD_400, align=PP_ALIGN.CENTER)

txt(sl, "Demanda = Mercado Objetivo × % Penetración   ·   Ingreso = Demanda × S/. 535.08 / año",
    0.45, 7.08, 12.43, 0.28, size=9.5, color=GRAY_500, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 7 — PROYECCIÓN 5 AÑOS
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Proyección 2027–2031", "Escenario Moderado  ·  Crecimiento anual: +15 %  ·  HolonIQ EdTech LATAM 2024")

# KPIs
kpis = [
    ("41,203 → 72,063", "Suscripciones  2027→2031", NAVY),
    ("+75 %",           "Crecimiento acumulado",      PURPLE_D),
    ("S/. 148.7M",      "Ingresos acumulados 5 años", GOLD_600),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.45 + i * 4.3
    _rect(sl, x, 1.52, 4.08, 1.22, fill=col)
    _rect(sl, x, 1.52, 4.08, 0.06, fill=GOLD_400)
    txt(sl, val, x+0.2, 1.6, 3.7, 0.62, size=20, bold=True, color=GOLD_400)
    txt(sl, lbl, x+0.2, 2.2, 3.7, 0.42, size=10, color=PURPLE_L)

# Barras
PROJ_YEARS = [2027, 2028, 2029, 2030, 2031]
PROJ_SUBS  = [41203, 47383, 54490, 62664, 72063]
PROJ_ING   = [22046259, 25353198, 29156177, 33529604, 38559045]
bar_cols   = [NAVY, PURPLE_D, PURPLE_M, GOLD_600, GOLD_500]

max_ing   = max(PROJ_ING)
chart_b   = 6.62
chart_h   = 3.3
bw        = 1.72
gap       = 0.48
sx        = 0.8

_rect(sl, sx-0.1, chart_b, 12.1, 0.05, fill=GOLD_500)

for i, (yr, sub, ing, col) in enumerate(zip(PROJ_YEARS, PROJ_SUBS, PROJ_ING, bar_cols)):
    x   = sx + i*(bw+gap)
    bh  = (ing/max_ing)*chart_h
    yb  = chart_b - bh
    _rect(sl, x, yb, bw, bh, fill=col)
    _rect(sl, x, yb, bw, 0.05, fill=GOLD_400)
    ing_m = ing/1_000_000
    txt(sl, f"S/.{ing_m:.1f}M", x-0.05, yb-0.42, bw+0.1, 0.35,
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, str(yr), x, chart_b+0.06, bw, 0.35,
        size=12, bold=True, color=GRAY_800, align=PP_ALIGN.CENTER)
    txt(sl, f"{sub:,}", x, chart_b+0.42, bw, 0.28,
        size=9, color=GRAY_500, align=PP_ALIGN.CENTER)

txt(sl, "Fuente: HolonIQ EdTech LATAM Report 2024  ·  CAGR 15 % para plataformas SaaS educativas en Perú y LATAM",
    0.45, 7.1, 12.43, 0.26, size=8, color=GRAY_500, italic=True)


# ════════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSIONES
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_content(sl)
epic_header(sl, "Conclusiones", "Síntesis del análisis de demanda — LegendaryClass")

concl = [
    ("82,406 docentes",   "conforman el mercado objetivo en Lima Metropolitana (MINEDU 2024).",     NAVY),
    ("S/. 535 / año",     "precio ponderado validado por disposición de pago real (Q17).",           PURPLE_D),
    ("S/. 22M en 2027",   "ingresos proyectados en escenario moderado (50 % penetración).",          GOLD_600),
    ("S/. 38.5M en 2031", "con crecimiento anual del 15 % — CAGR EdTech LATAM (HolonIQ 2024).",    PURPLE_M),
    ("58 % de docentes",  "reportan baja motivación frecuente — necesidad real y cuantificada.",     GRAY_700),
]
for i, (val, rest, col) in enumerate(concl):
    y = 1.52 + i * 1.03
    card(sl, 0.45, y, 12.43, 0.9)
    _rect(sl, 0.45, y, 0.12, 0.9, fill=col)
    _rect(sl, 0.45, y, 12.43, 0.05, fill=GOLD_500)
    txt(sl, str(i+1), 0.65, y+0.2, 0.48, 0.5,
        size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, val,  1.28, y+0.17, 3.3, 0.55, size=17, bold=True, color=col)
    txt(sl, rest, 4.65, y+0.2,  7.9, 0.52, size=12, color=GRAY_700)


# ════════════════════════════════════════════════════════
# SLIDE 9 — CIERRE
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
if os.path.exists(FONDO):
    pic = sl.shapes.add_picture(FONDO, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
overlay2 = _rect(sl, 0, 0, 13.33, 7.5, fill=RGB(255,255,255))
sp2 = overlay2._element
spPr2 = sp2.find(qn('p:spPr'))
sf2 = spPr2.find('.//' + qn('a:solidFill'))
if sf2 is not None:
    sg2 = sf2.find(qn('a:srgbClr'))
    if sg2 is None: sg2 = etree.SubElement(sf2, qn('a:srgbClr'))
    sg2.set('val', 'FFFFFF')
    al2 = etree.SubElement(sg2, qn('a:alpha'))
    al2.set('val', '82000')

_rect(sl, 0, 0, 13.33, 1.6, fill=NAVY)
_rect(sl, 2.0, 0, 9.33, 1.6, fill=PURPLE_D)
_rect(sl, 4.5, 0, 4.33, 1.6, fill=PURPLE_M)
_rect(sl, 0, 1.55, 13.33, 0.07, fill=GOLD_500)

txt(sl, "LEGENDARYCLASS", 0.5, 0.18, 12.33, 1.0,
    size=52, bold=True, color=GOLD_400, align=PP_ALIGN.CENTER)

card(sl, 3.5, 2.2, 6.33, 3.3)
txt(sl, "⚔️", 5.9, 2.35, 1.5, 0.65, size=30, align=PP_ALIGN.CENTER)
txt(sl, "Sistema SaaS de Gamificación Educativa",
    3.7, 3.05, 5.93, 0.45, size=12, color=GRAY_500, align=PP_ALIGN.CENTER, italic=True)
gold_line(sl, 4.1, 3.6, 5.13)
txt(sl, "¡Gracias!", 3.7, 3.75, 5.93, 0.75,
    size=34, bold=True, color=GOLD_500, align=PP_ALIGN.CENTER)
gold_line(sl, 4.1, 4.58, 5.13)
txt(sl, f"{ALUMNO}  ·  {CURSO}  ·  {ANIO}",
    3.7, 4.7, 5.93, 0.38, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)
txt(sl, f"{INSTITUCION}  ·  {CARRERA}",
    3.7, 5.1, 5.93, 0.38, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)

# ── Guardar ──────────────────────────────────────────────
OUT = "C:/Proyectos/LegendaryClass/lab03/Presentacion_Caso3_LegendaryClass_v4.pptx"
prs.save(OUT)
print("[OK] PPT guardada:", OUT)
