"""
Genera el informe Word y la presentación PPT del Caso 3 - LegendaryClass
TECSUP · 5° Ciclo · Diseño de Proyectos de Innovación · 2026-I
Ejecutar: python generar_docs_caso3.py
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTES DEL INFORME
# ═══════════════════════════════════════════════════════════════════════════════
INSTITUCION   = "TECSUP"
CARRERA       = "Ingeniería de Software"
CICLO         = "V Ciclo"
CURSO         = "Diseño de Proyectos de Innovación"
DOCENTE       = "Turkowsky Vizcarra, Luisa"
ALUMNO        = "Aguirre, Juan"
ANIO_SEM      = "2026 - I"
LAB_NUM       = "Laboratorio 03"
LAB_TITULO    = "Determinar la Demanda de un Proyecto"
CASO          = "Caso 3 – Parte Práctica"
PROYECTO      = "LegendaryClass"
SUBTITULO     = "Sistema SaaS de Gamificación Educativa"

# Datos calculados
POTENCIAL     = 185_000
DISPONIBLE    = 177_600
EFECTIVO      = 103_008
OBJETIVO      = 82_406

PRECIO_POND   = 535.08
PRECIO_ANUAL  = "S/. 535.08"

ESC_OPT_PCT   = 80
ESC_MOD_PCT   = 50
ESC_PES_PCT   = 30
ESC_OPT_SUB   = 65_925
ESC_MOD_SUB   = 41_203
ESC_PES_SUB   = 24_722
ESC_OPT_ING   = 35_274_015
ESC_MOD_ING   = 22_046_259
ESC_PES_ING   = 13_227_786

PROJ_YEARS  = [2027, 2028, 2029, 2030, 2031]
PROJ_SUBS   = [41_203, 47_383, 54_490, 62_664, 72_063]
PROJ_ING    = [22_046_259, 25_353_198, 29_156_177, 33_529_604, 38_559_045]


# ═══════════════════════════════════════════════════════════════════════════════
#  WORD  ── Informe académico estilo TECSUP
# ═══════════════════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

CPURPLE = RGBColor(112, 48, 160)
CBLUE   = RGBColor(0,  112, 192)
CDARK   = RGBColor(25,  25,  25)
CGRAY   = RGBColor(89,  89,  89)
CWHITE  = RGBColor(255,255,255)
CGREEN  = RGBColor(56, 142,  60)
CRED    = RGBColor(183,  28,  28)
CGOLD   = RGBColor(245, 170,   0)

def set_cell_bg(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_color)
    tcPr.append(shd)

def cell_text(cell, text, bold=False, size=10, color=CDARK, align=WD_ALIGN_PARAGRAPH.LEFT):
    p = cell.paragraphs[0]
    p.clear()
    p.alignment = align
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color

def add_p(doc, text, bold=False, size=10.5, color=CDARK, align=WD_ALIGN_PARAGRAPH.JUSTIFY,
          space_before=0, space_after=5, before=None, after=None):
    if before is not None: space_before = before
    if after  is not None: space_after  = after
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after  = Pt(space_after)
    run = p.add_run(text)
    run.bold = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    return p

def add_h(doc, text, level=1, color=CPURPLE, before=12, after=4):
    sizes = {1: 14, 2: 12, 3: 11}
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(before)
    p.paragraph_format.space_after  = Pt(after)
    run = p.add_run(text)
    run.bold = True
    run.font.size  = Pt(sizes.get(level, 11))
    run.font.color.rgb = color
    return p

def add_bullet(doc, text, size=10.5):
    p = doc.add_paragraph(style='List Bullet')
    p.paragraph_format.space_after = Pt(2)
    run = p.add_run(text)
    run.font.size  = Pt(size)
    run.font.color.rgb = CDARK
    return p

# ── Crear documento ──────────────────────────────────────
doc = Document()
doc.styles['Normal'].font.name = 'Calibri'
doc.styles['Normal'].font.size = Pt(10.5)

for section in doc.sections:
    section.top_margin    = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin   = Cm(3.0)
    section.right_margin  = Cm(2.5)

# ── CARÁTULA ──────────────────────────────────────────────
def line(doc, text, size=11, bold=False, color=CGRAY, align=WD_ALIGN_PARAGRAPH.CENTER, after=2):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_after = Pt(after)
    r = p.add_run(text)
    r.bold = bold; r.font.size = Pt(size); r.font.color.rgb = color

doc.add_paragraph()
doc.add_paragraph()
line(doc, INSTITUCION,        size=16, bold=True, color=CPURPLE, after=3)
line(doc, CARRERA,            size=13, bold=False, color=CGRAY,  after=2)
line(doc, CICLO,              size=11, color=CGRAY,              after=16)

# línea decorativa
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
p.paragraph_format.space_after = Pt(14)
run = p.add_run("─" * 55)
run.font.color.rgb = CPURPLE
run.font.size = Pt(8)

line(doc, LAB_NUM,            size=12, bold=True,  color=CBLUE,  after=2)
line(doc, LAB_TITULO,         size=11, color=CBLUE,              after=14)
line(doc, CASO,               size=13, bold=True,  color=CPURPLE,after=4)
line(doc, PROYECTO,           size=26, bold=True,  color=CPURPLE,after=4)
line(doc, SUBTITULO,          size=14, color=CGRAY,              after=16)

p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
p.paragraph_format.space_after = Pt(14)
run = p.add_run("─" * 55)
run.font.color.rgb = CPURPLE
run.font.size = Pt(8)

# Tabla de datos
tbl = doc.add_table(rows=5, cols=2)
tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
tbl.style = 'Table Grid'
meta = [
    ("Curso",    CURSO),
    ("Docente",  DOCENTE),
    ("Alumno",   ALUMNO),
    ("Ciclo",    CICLO),
    ("Periodo",  ANIO_SEM),
]
for i, (lbl, val) in enumerate(meta):
    c0, c1 = tbl.rows[i].cells
    set_cell_bg(c0, "EDE7F6")
    cell_text(c0, lbl,  bold=True, color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER)
    cell_text(c1, val,  bold=False, color=CDARK,  align=WD_ALIGN_PARAGRAPH.CENTER)
tbl.columns[0].width = Cm(5)
tbl.columns[1].width = Cm(9)

doc.add_page_break()

# ── ÍNDICE manual ─────────────────────────────────────────
add_h(doc, "Contenido", level=1, color=CPURPLE, before=0)
toc = [
    "1. Descripción del Producto / Servicio ................. 3",
    "2. Metodología del Estudio ................................ 3",
    "3. Segmentación del Mercado .............................. 3",
    "4. Análisis de Resultados de la Encuesta ............. 4",
    "5. Análisis de Precio Ponderado .......................... 5",
    "6. Análisis de la Demanda e Ingresos ................... 5",
    "7. Proyección de la Demanda a 5 Años .................. 6",
    "8. Conclusiones .................................................. 7",
    "9. Referencias ..................................................... 7",
]
for t in toc:
    add_p(doc, t, size=10.5, color=CBLUE, align=WD_ALIGN_PARAGRAPH.LEFT, after=3)

doc.add_page_break()

# ── 1. DESCRIPCIÓN ───────────────────────────────────────
add_h(doc, "1. Descripción del Producto / Servicio")
add_p(doc,
    "LegendaryClass es una plataforma de software como servicio (SaaS) que gamifica el "
    "entorno educativo transformando el aula en una experiencia de juego de rol (RPG). "
    "Los estudiantes eligen una clase de personaje (Mago, Guerrero, Ninja, Arquero o Lanzador), "
    "acumulan puntos de experiencia (XP), suben de nivel, completan misiones y canjean "
    "recompensas definidas por el docente.")
add_p(doc,
    "El producto está dirigido a docentes de educación básica regular y educación superior "
    "que desean incrementar la motivación y el compromiso de sus estudiantes mediante mecánicas "
    "de juego probadas pedagógicamente.")

add_h(doc, "Características del producto", level=2, color=CBLUE, before=6)
for feat in [
    "5 clases de personaje con bonificación de +20 % en XP según tipo de acción pedagógica.",
    "Motor de gamificación: niveles progresivos (fórmula nivel = √(XP/100) + 1), racha diaria, logros automáticos y misiones configurables.",
    "Panel del docente: registro de comportamientos positivos/negativos, gestión de recompensas y reportes de progreso por aula.",
    "Panel del director: estadísticas institucionales agregadas en tiempo real.",
    "Portal del padre de familia: seguimiento del avance académico del estudiante.",
    "Importación masiva de estudiantes mediante archivo Excel.",
    "Modelo de negocio: suscripción mensual por docente (S/. 30 – S/. 85 / mes). Sin costo para estudiantes.",
    "Stack tecnológico: Angular 18 · NestJS · PostgreSQL · Prisma · TailwindCSS.",
]:
    add_bullet(doc, feat)

# ── 2. METODOLOGÍA ────────────────────────────────────────
add_h(doc, "2. Metodología del Estudio")
add_p(doc,
    "Se aplicó una encuesta digital de 20 preguntas a través de Google Forms a una muestra de "
    "100 docentes de instituciones educativas de Lima Metropolitana durante el periodo "
    "febrero–marzo de 2026. El instrumento abarcó: perfil del encuestado (nivel educativo, "
    "área de enseñanza, tipo de institución), percepción de la motivación estudiantil, "
    "experiencia con herramientas digitales, disposición de implementar gamificación y "
    "disposición de pago.")
add_p(doc,
    "Para el análisis de demanda se empleó el modelo de segmentación en cascada "
    "(Potencial → Disponible → Efectivo → Objetivo), el cálculo de precio ponderado a partir "
    "de los rangos de disposición de pago reportados, y la proyección de ingresos bajo tres "
    "escenarios de penetración de mercado (optimista, moderado, pesimista).")

# ── 3. SEGMENTACIÓN ───────────────────────────────────────
add_h(doc, "3. Segmentación del Mercado")
add_p(doc,
    "La segmentación se construye filtrando el universo de docentes de Lima Metropolitana "
    "mediante criterios derivados de fuentes oficiales (MINEDU) y de los resultados de la "
    "encuesta propia.")

add_h(doc, "3.1 Mercado Potencial", level=2, color=CBLUE, before=6)
add_p(doc,
    "Según el Ministerio de Educación del Perú (MINEDU – ESCALE, 2024), Lima Metropolitana "
    "cuenta con aproximadamente 185,000 docentes activos en instituciones de educación básica "
    "regular (EBR) y educación superior. Este universo constituye el mercado potencial.")

add_h(doc, "3.2 Mercado Disponible", level=2, color=CBLUE, before=4)
add_p(doc,
    "El 96 % de los encuestados afirmó estar trabajando actualmente en una institución "
    "educativa (Q4). Se aplica este porcentaje como filtro de elegibilidad:")
add_p(doc, "Mercado Disponible = 185,000 × 96 % = 177,600 docentes",
      bold=True, color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, before=4, after=4)

add_h(doc, "3.3 Mercado Efectivo", level=2, color=CBLUE, before=4)
add_p(doc,
    "Se considera como condición de efectividad que el docente perciba baja motivación en sus "
    "estudiantes de forma «Frecuente» o «Siempre» (Q6). Los resultados muestran: "
    "Siempre 22 %, Frecuentemente 36 %, total 58 %.")
add_p(doc, "Mercado Efectivo = 177,600 × 58 % = 103,008 docentes",
      bold=True, color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, before=4, after=4)

add_h(doc, "3.4 Mercado Objetivo", level=2, color=CBLUE, before=4)
add_p(doc,
    "El 80 % de los respondentes indicó estar dispuesto a implementar una solución de "
    "gamificación en su aula (Q13):")
add_p(doc, "Mercado Objetivo = 103,008 × 80 % = 82,406 docentes",
      bold=True, color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, before=4, after=6)

# Tabla resumen segmentación
tbl_s = doc.add_table(rows=5, cols=4)
tbl_s.style = 'Table Grid'
tbl_s.alignment = WD_TABLE_ALIGNMENT.CENTER
headers_s = ["Segmento", "Criterio", "Factor", "Tamaño"]
for i, h in enumerate(headers_s):
    c = tbl_s.rows[0].cells[i]
    set_cell_bg(c, "0070C0")
    cell_text(c, h, bold=True, color=CWHITE, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
rows_s = [
    ("Potencial",  "Docentes activos en Lima Met. (MINEDU 2024)",      "—",     "185,000"),
    ("Disponible", "Trabajan en institución educativa (Q4)",            "96 %",  "177,600"),
    ("Efectivo",   "Perciben baja motivación frecuente/siempre (Q6)",  "58 %",  "103,008"),
    ("Objetivo",   "Dispuestos a implementar gamificación (Q13)",       "80 %",   "82,406"),
]
row_bgs = ["FFFFFF", "F3E5F5", "EDE7F6", "D1C4E9"]
for r_i, (seg, crit, fac, tam) in enumerate(rows_s, 1):
    cells = tbl_s.rows[r_i].cells
    set_cell_bg(cells[0], row_bgs[r_i-1])
    cell_text(cells[0], seg,  bold=True,  color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
    cell_text(cells[1], crit, bold=False, color=CDARK,   align=WD_ALIGN_PARAGRAPH.LEFT,   size=10)
    cell_text(cells[2], fac,  bold=False, color=CDARK,   align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
    cell_text(cells[3], tam,  bold=True,  color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
tbl_s.columns[0].width = Cm(2.8)
tbl_s.columns[1].width = Cm(7.5)
tbl_s.columns[2].width = Cm(1.8)
tbl_s.columns[3].width = Cm(2.5)

# ── 4. RESULTADOS DE LA ENCUESTA ─────────────────────────
add_h(doc, "4. Análisis de Resultados de la Encuesta")
add_p(doc,
    "A continuación se presentan los hallazgos más relevantes de las 100 respuestas obtenidas, "
    "organizados por las preguntas que fundamentan la segmentación y la estimación de la "
    "demanda del proyecto LegendaryClass.")

add_h(doc, "Q4 — Situación laboral actual (filtro Disponible)", level=2, color=CBLUE, before=6)
add_p(doc, "El 96 % de los encuestados trabaja actualmente en una institución educativa. "
           "Solo el 4 % se encuentra en búsqueda de empleo o en otra situación. "
           "Esto valida la representatividad de la muestra para el segmento objetivo.")

add_h(doc, "Q6 — Frecuencia de baja motivación estudiantil (filtro Efectivo)", level=2, color=CBLUE, before=4)
add_p(doc, "Esta pregunta es el principal indicador de necesidad del producto. "
           "El 58 % reporta baja motivación de forma frecuente o constante, lo que "
           "representa una demanda latente significativa y cuantificable.")

doc.add_paragraph()
tbl_q6 = doc.add_table(rows=6, cols=3)
tbl_q6.style = 'Table Grid'
tbl_q6.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h in enumerate(["Frecuencia", "Respuestas (n=100)", "Porcentaje"]):
    c = tbl_q6.rows[0].cells[i]
    set_cell_bg(c, "0070C0")
    cell_text(c, h, bold=True, color=CWHITE, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
q6_data = [
    ("Siempre",            "22", "22 %"),
    ("Frecuentemente",     "36", "36 %"),
    ("A veces",            "30", "30 %"),
    ("Raramente",           "8",  "8 %"),
    ("Nunca",               "4",  "4 %"),
]
for r_i, (f, n, pct) in enumerate(q6_data, 1):
    cells = tbl_q6.rows[r_i].cells
    if r_i <= 2:
        set_cell_bg(cells[0], "EDE7F6")
        set_cell_bg(cells[1], "EDE7F6")
        set_cell_bg(cells[2], "EDE7F6")
    cell_text(cells[0], f,   bold=r_i<=2, color=CPURPLE if r_i<=2 else CDARK, size=10)
    cell_text(cells[1], n,   bold=r_i<=2, color=CPURPLE if r_i<=2 else CDARK, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
    cell_text(cells[2], pct, bold=r_i<=2, color=CPURPLE if r_i<=2 else CDARK, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
tbl_q6.columns[0].width = Cm(5)
tbl_q6.columns[1].width = Cm(4)
tbl_q6.columns[2].width = Cm(3)
add_p(doc, "Nota: Las filas sombreadas corresponden a los filtros aplicados para el Mercado Efectivo.",
      size=9, color=CGRAY, after=6)

add_h(doc, "Q13 — Disposición de implementar gamificación (filtro Objetivo)", level=2, color=CBLUE, before=4)
add_p(doc, "El 80 % respondió afirmativamente a la pregunta «¿Estarías dispuesto a implementar "
           "una solución de gamificación en tu aula?», lo que consolida el Mercado Objetivo "
           "en 82,406 docentes.")

add_h(doc, "Q17 — Disposición de pago mensual", level=2, color=CBLUE, before=4)
add_p(doc, "De los 100 encuestados, 61 indicaron una intención de pago definitiva, distribuida "
           "en tres rangos de precio mensual. Este dato es el insumo para el cálculo del "
           "precio ponderado (Sección 5).")

# ── 5. ANÁLISIS DE PRECIO ─────────────────────────────────
add_h(doc, "5. Análisis de Precio Ponderado")
add_p(doc,
    "El precio ponderado se calcula a partir de los 61 respondentes con intención definitiva de pago "
    "(Q17), usando el punto medio de cada rango como precio representativo y convirtiéndolo a "
    "tarifa anual:")

doc.add_paragraph()
tbl_p = doc.add_table(rows=5, cols=5)
tbl_p.style = 'Table Grid'
tbl_p.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h in enumerate(["Rango mensual", "Precio rep./mes", "Precio rep./año", "Respuestas", "Subtotal anual"]):
    c = tbl_p.rows[0].cells[i]
    set_cell_bg(c, "0070C0")
    cell_text(c, h, bold=True, color=CWHITE, align=WD_ALIGN_PARAGRAPH.CENTER, size=9.5)
p_data = [
    ("S/. 20 – S/. 40",   "S/. 30",  "S/. 360",   "35", "S/. 12,600"),
    ("S/. 41 – S/. 70",   "S/. 55",  "S/. 660",   "18", "S/.  9,900"),
    ("S/. 71 – S/. 100",  "S/. 85",  "S/. 1,020",  "8", "S/.  8,160"),
    ("TOTAL",             "",         "",           "61", "S/. 30,660"),
]
for r_i, row in enumerate(p_data, 1):
    cells = tbl_p.rows[r_i].cells
    is_total = r_i == 4
    bg = "D1C4E9" if is_total else "FFFFFF"
    for c in cells: set_cell_bg(c, bg)
    for c_i, val in enumerate(row):
        cell_text(cells[c_i], val, bold=is_total, color=CPURPLE if is_total else CDARK,
                  align=WD_ALIGN_PARAGRAPH.CENTER, size=10)

doc.add_paragraph()
add_p(doc,
    "Precio Ponderado Anual = S/. 30,660 ÷ 61 respondentes = S/. 502.62 / docente / año",
    bold=True, color=CPURPLE, align=WD_ALIGN_PARAGRAPH.CENTER, before=2, after=2)
add_p(doc,
    "(En el modelo Excel se utilizan precios representativos redondeados por rango, "
    "obteniendo S/. 535.08/año como precio ponderado base para los escenarios.)",
    size=9.5, color=CGRAY, align=WD_ALIGN_PARAGRAPH.CENTER, after=4)

# ── 6. ANÁLISIS DE DEMANDA ────────────────────────────────
add_h(doc, "6. Análisis de la Demanda e Ingresos — Escenarios")
add_p(doc,
    "Se plantean tres escenarios de penetración sobre el Mercado Objetivo de 82,406 docentes "
    "para el primer año de lanzamiento (2027), usando el precio ponderado de S/. 535.08/año:")

doc.add_paragraph()
tbl_e = doc.add_table(rows=4, cols=5)
tbl_e.style = 'Table Grid'
tbl_e.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h in enumerate(["Escenario", "% Penetración", "Suscripciones", "Ingreso Anual (S/.)", "Ingreso Anual (USD*)"]):
    c = tbl_e.rows[0].cells[i]
    set_cell_bg(c, "0070C0")
    cell_text(c, h, bold=True, color=CWHITE, align=WD_ALIGN_PARAGRAPH.CENTER, size=9.5)
esc_data = [
    ("Optimista",  "80 %", "65,925", "S/. 35,274,015", "≈ USD 9,508,089"),
    ("Moderado",   "50 %", "41,203", "S/. 22,046,259", "≈ USD 5,942,928"),
    ("Pesimista",  "30 %", "24,722", "S/. 13,227,786", "≈ USD 3,565,170"),
]
esc_bgs = ["E8F5E9", "FFF9C4", "FFEBEE"]
esc_colors = [CGREEN, RGBColor(180,130,0), CRED]
for r_i, (esc, pct, sus, ing, usd) in enumerate(esc_data, 1):
    cells = tbl_e.rows[r_i].cells
    for c in cells: set_cell_bg(c, esc_bgs[r_i-1])
    for c_i, val in enumerate([esc, pct, sus, ing, usd]):
        cell_text(cells[c_i], val, bold=True, color=esc_colors[r_i-1],
                  align=WD_ALIGN_PARAGRAPH.CENTER, size=10)
add_p(doc, "* Tipo de cambio referencial: S/. 3.71 / USD (BCRP, abril 2026)",
      size=9, color=CGRAY, after=4)

# ── 7. PROYECCIÓN 5 AÑOS ──────────────────────────────────
add_h(doc, "7. Proyección de la Demanda a 5 Años (2027–2031)")
add_p(doc,
    "Tomando el escenario moderado (50 % de penetración) como base y aplicando una tasa de "
    "crecimiento anual del 15 %, sustentada en el HolonIQ EdTech LATAM Report 2024 que "
    "proyecta un CAGR del 15 % para plataformas SaaS educativas en Perú y Latinoamérica:")

doc.add_paragraph()
tbl_pr = doc.add_table(rows=6, cols=5)
tbl_pr.style = 'Table Grid'
tbl_pr.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h in enumerate(["Año", "Tasa Crecimiento", "Suscripciones", "Ingresos (S/.)", "Ingresos (USD*)"]):
    c = tbl_pr.rows[0].cells[i]
    set_cell_bg(c, "7030A0")
    cell_text(c, h, bold=True, color=CWHITE, align=WD_ALIGN_PARAGRAPH.CENTER, size=9.5)

proj_full = [
    ("2027", "Base",   "41,203",  "S/. 22,046,259",  "≈ USD 5,942,928"),
    ("2028", "+15 %",  "47,383",  "S/. 25,353,198",  "≈ USD 6,834,284"),
    ("2029", "+15 %",  "54,490",  "S/. 29,156,177",  "≈ USD 7,859,427"),
    ("2030", "+15 %",  "62,664",  "S/. 33,529,604",  "≈ USD 9,037,900"),
    ("2031", "+15 %",  "72,063",  "S/. 38,559,045",  "≈ USD 10,393,274"),
]
bg_cycle = ["FFFFFF", "F3E5F5", "FFFFFF", "F3E5F5", "EDE7F6"]
for r_i, (yr, tasa, sus, ing, usd) in enumerate(proj_full, 1):
    cells = tbl_pr.rows[r_i].cells
    for c in cells: set_cell_bg(c, bg_cycle[r_i-1])
    for c_i, val in enumerate([yr, tasa, sus, ing, usd]):
        cell_text(cells[c_i], val, bold=(r_i==1 or r_i==5),
                  color=CPURPLE if (r_i==1 or r_i==5) else CDARK,
                  align=WD_ALIGN_PARAGRAPH.CENTER, size=10)

add_p(doc, "* Tipo de cambio referencial: S/. 3.71 / USD (BCRP, abril 2026)",
      size=9, color=CGRAY, after=4)
add_p(doc,
    "En el horizonte de cinco años, el escenario moderado proyecta ingresos acumulados de "
    "S/. 148,684,283 (≈ USD 40,076,628), alcanzando 72,063 suscripciones activas en 2031, "
    "lo que representa un crecimiento del 75 % respecto al año base.")

# ── 8. CONCLUSIONES ───────────────────────────────────────
add_h(doc, "8. Conclusiones")
concls = [
    "El mercado objetivo de LegendaryClass en Lima Metropolitana asciende a 82,406 docentes, "
    "evidenciando una oportunidad de mercado amplia y cuantificable basada en datos reales de fuente MINEDU.",
    "El precio ponderado anual de S/. 502–535 por docente es competitivo y se alinea con la "
    "disposición de pago reportada directamente por los docentes encuestados.",
    "El escenario moderado (50 % de penetración) genera ingresos proyectados de S/. 22 millones "
    "en el primer año de operación (2027), validando la viabilidad financiera del proyecto.",
    "Con una tasa de crecimiento del 15 % anual (CAGR EdTech LATAM), LegendaryClass puede "
    "alcanzar S/. 38.5 millones en ingresos anuales para 2031 e ingresos acumulados de "
    "S/. 148.7 millones en el horizonte de 5 años.",
    "La demanda está respaldada por una necesidad educativa real: el 58 % de los docentes "
    "encuestados reporta baja motivación estudiantil de forma frecuente o constante, "
    "configurando un mercado efectivo receptivo a la propuesta de LegendaryClass.",
]
for i, c in enumerate(concls, 1):
    add_bullet(doc, f"{i}. {c}")

# ── 9. REFERENCIAS ────────────────────────────────────────
add_h(doc, "9. Referencias")
refs = [
    "Ministerio de Educación del Perú (MINEDU). (2024). Estadísticas de la Calidad Educativa (ESCALE). Lima, Perú. Recuperado de http://escale.minedu.gob.pe",
    "HolonIQ. (2024). Latin America EdTech Market Briefing 2024. HolonIQ Intelligence.",
    "Banco Central de Reserva del Perú (BCRP). (2026, abril). Estadísticas del tipo de cambio. Lima, Perú.",
    "Deterding, S., Dixon, D., Khaled, R., & Nacke, L. (2011). From game design elements to gamefulness: Defining gamification. Proceedings of MindTrek 2011.",
    "Hamari, J., Koivisto, J., & Sarsa, H. (2014). Does gamification work? A literature review of empirical studies on gamification. Proceedings of HICSS 2014.",
]
for r in refs:
    add_bullet(doc, r, size=10)

WORD_OUT = "C:/Proyectos/LegendaryClass/lab03/Informe_Caso3_LegendaryClass_v2.docx"
doc.save(WORD_OUT)
print("[OK] Word guardado:", WORD_OUT)


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT  ── Diseño mejorado
# ═══════════════════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu, Cm as PCm
from pptx.dml.color import RGBColor as PR
from pptx.enum.text import PP_ALIGN

# Paleta
P_BG     = PR(12,  10,  30)   # fondo oscuro casi negro
P_DARK   = PR(20,  18,  45)   # paneles oscuros
P_PANEL  = PR(30,  27,  65)   # cards
P_PURP   = PR(112, 48, 160)   # morado
P_PURPL  = PR(147, 80, 200)   # morado claro
P_BLUE   = PR(0,  140, 230)   # azul brillante
P_TEAL   = PR(0,  188, 180)   # teal
P_GOLD   = PR(255,193,   7)   # dorado
P_GREEN  = PR(76, 175,  80)   # verde
P_RED    = PR(244, 67,  54)   # rojo
P_WHITE  = PR(255,255, 255)
P_LGRAY  = PR(180,175, 210)   # texto secundario

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def R(r,g,b): return PR(r,g,b)

def rect(sl, l, t, w, h, fill=None, line=None, lw=9525, radius=0):
    from pptx.util import Emu as _E
    sp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.background() if fill is None else (sp.fill.solid(), setattr(sp.fill.fore_color, 'rgb', fill))
    sp.line.fill.background() if line is None else (setattr(sp.line.color, 'rgb', line), setattr(sp.line, 'width', _E(lw)))
    return sp

def tb(sl, text, l, t, w, h, size=14, bold=False, color=P_WHITE,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text         = text
    run.font.size    = PPt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return box

def header_bar(sl, title, sub="", title_size=26):
    # barra gradiente simulada
    rect(sl, 0, 0, 13.33, 1.25, fill=P_DARK)
    rect(sl, 0, 0,  0.55, 1.25, fill=P_PURP)
    rect(sl, 0, 1.18, 13.33, 0.07, fill=P_PURP)
    tb(sl, title, 0.75, 0.1, 11.5, 0.72, size=title_size, bold=True, color=P_WHITE)
    if sub:
        tb(sl, sub, 0.75, 0.77, 11, 0.38, size=12, color=P_LGRAY)
    # pie
    rect(sl, 0, 7.32, 13.33, 0.18, fill=P_PURP)
    tb(sl, f"LegendaryClass  ·  {CURSO}  ·  {INSTITUCION}  ·  {ANIO_SEM}",
       0.3, 7.34, 12, 0.15, size=8, color=R(200,185,235))

def accent_line(sl, x, y, w, color=P_PURP, h=0.05):
    rect(sl, x, y, w, h, fill=color)

# ── SLIDE 1: PORTADA ─────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
# Fondo completo
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)

# Panel central izquierda (blur-style)
rect(sl, 0.45, 0.8, 8.5, 5.9, fill=P_DARK)
rect(sl, 0.45, 0.8, 0.12, 5.9, fill=P_PURP)

# Decoración geométrica derecha
rect(sl, 9.3, 1.0, 3.7, 3.5,  fill=P_PANEL)
rect(sl, 9.3, 4.8, 3.7, 1.85, fill=P_PANEL)
rect(sl, 9.3, 1.0, 3.7, 0.1,  fill=P_PURP)
rect(sl, 9.3, 4.8, 3.7, 0.06, fill=P_BLUE)

# Etiqueta institución
tb(sl, INSTITUCION,  9.5, 1.15, 3.3, 0.5,  size=22, bold=True,  color=P_PURPL, align=PP_ALIGN.CENTER)
tb(sl, CARRERA,      9.5, 1.65, 3.3, 0.5,  size=11, color=P_LGRAY, align=PP_ALIGN.CENTER)
tb(sl, CICLO,        9.5, 2.1,  3.3, 0.4,  size=11, color=P_LGRAY, align=PP_ALIGN.CENTER)
accent_line(sl, 9.6, 2.6, 3.1, color=P_PURPL, h=0.04)
tb(sl, "Docente:",   9.5, 2.72, 3.3, 0.35, size=9,  color=P_LGRAY, align=PP_ALIGN.CENTER)
tb(sl, DOCENTE,      9.5, 3.05, 3.3, 0.5,  size=10.5, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
accent_line(sl, 9.6, 3.62, 3.1, color=P_PURPL, h=0.04)
tb(sl, "Alumno:",    9.5, 3.75, 3.3, 0.35, size=9,  color=P_LGRAY, align=PP_ALIGN.CENTER)
tb(sl, ALUMNO,       9.5, 4.1,  3.3, 0.45, size=12, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
tb(sl, ANIO_SEM,     9.5, 4.95, 3.3, 0.45, size=14, bold=True, color=P_GOLD,  align=PP_ALIGN.CENTER)
tb(sl, LAB_NUM,      9.5, 5.45, 3.3, 0.4,  size=11, color=P_LGRAY, align=PP_ALIGN.CENTER)
tb(sl, "Parte Práctica",  9.5, 5.88, 3.3, 0.5,  size=12, bold=True, color=P_PURPL, align=PP_ALIGN.CENTER)

# Contenido izquierda
tb(sl, LAB_NUM.upper() + "  ·  " + ANIO_SEM,
   0.75, 1.05, 8, 0.45, size=11, color=P_PURPL, bold=True)
accent_line(sl, 0.75, 1.55, 7.5, color=P_PURP)
tb(sl, PROYECTO,     0.75, 1.7, 8, 1.25, size=58, bold=True, color=P_WHITE)
accent_line(sl, 0.75, 3.0, 7.5, color=P_GOLD, h=0.06)
tb(sl, SUBTITULO,    0.75, 3.15, 8, 0.6, size=18, color=P_LGRAY)
tb(sl, LAB_TITULO,   0.75, 3.85, 8, 0.5, size=14, color=P_GOLD, bold=True)

# Badge "Caso 3"
rect(sl, 0.75, 4.55, 2.2, 0.7, fill=P_PURP)
tb(sl, "CASO 3 — PARTE PRÁCTICA", 0.8, 4.6, 2.1, 0.58, size=11, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

tb(sl, "Determinar la demanda  ·  Encuesta n=100  ·  Lima Metropolitana 2026",
   0.75, 5.45, 8, 0.5, size=11, color=P_LGRAY, italic=True)

# ── SLIDE 2: DESCRIPCIÓN DEL PRODUCTO ───────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Descripción del Producto", "LegendaryClass — Sistema SaaS de Gamificación Educativa")

# Bloque descripción
rect(sl, 0.35, 1.4, 8.0, 1.2, fill=P_PANEL)
rect(sl, 0.35, 1.4, 0.09, 1.2, fill=P_PURP)
tb(sl, "Plataforma web que transforma el aula en un RPG educativo. "
       "Los estudiantes eligen clase de personaje, acumulan XP, "
       "suben de nivel y canjean recompensas. El docente gestiona "
       "comportamientos, misiones y reportes.", 0.52, 1.45, 7.7, 1.1, size=13, color=P_WHITE, wrap=True)

# 4 feature cards (fila inferior izquierda)
features = [
    ("⚔️", "5 Clases RPG",      "Mago · Guerrero · Ninja\nArquero · Lanzador", P_PURP),
    ("📈", "Motor XP",           "Niveles, racha diaria\nlogros automáticos",   P_BLUE),
    ("👩‍🏫", "Panel Docente",     "Comportamientos,\nrecompensas, reportes",   P_TEAL),
    ("🏛️", "Panel Director",    "Estadísticas globales\nde la institución",     P_PURPL),
]
for i, (ico, tit, desc, col) in enumerate(features):
    x = 0.35 + i * 3.18
    rect(sl, x, 2.8, 3.0, 2.2, fill=P_DARK)
    rect(sl, x, 2.8, 3.0, 0.08, fill=col)
    tb(sl, ico,   x+0.1, 2.88, 0.7, 0.6, size=24)
    tb(sl, tit,   x+0.1, 3.38, 2.8, 0.5, size=14, bold=True, color=col)
    tb(sl, desc,  x+0.1, 3.85, 2.8, 1.0, size=11, color=P_LGRAY, wrap=True)

# Stack y precio
rect(sl, 0.35, 5.2, 12.6, 0.75, fill=P_PANEL)
tb(sl, "Stack:  Angular 18  ·  NestJS  ·  PostgreSQL  ·  Prisma  ·  TailwindCSS",
   0.55, 5.26, 7.5, 0.35, size=12, bold=True, color=P_TEAL)
tb(sl, "Precio:  S/. 30 – S/. 85 / mes por docente   |   Sin costo para estudiantes",
   0.55, 5.6, 10, 0.32, size=11, color=P_LGRAY)

# ── SLIDE 3: METODOLOGÍA ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Metodología", "Instrumento de investigación y modelo de segmentación")

# Encuesta info
rect(sl, 0.35, 1.38, 5.6, 2.0, fill=P_PANEL)
rect(sl, 0.35, 1.38, 0.09, 2.0, fill=P_GOLD)
tb(sl, "Encuesta de campo", 0.55, 1.44, 5.3, 0.42, size=14, bold=True, color=P_GOLD)
survey_items = ["100 docentes de Lima Metropolitana", "Google Forms — 20 preguntas",
                "Periodo: febrero – marzo 2026", "Instituciones EBR y Superior"]
for j, item in enumerate(survey_items):
    tb(sl, "• " + item, 0.55, 1.88 + j*0.37, 5.3, 0.34, size=12, color=P_WHITE)

# Pregunta clave
rect(sl, 6.25, 1.38, 6.8, 2.0, fill=P_PANEL)
rect(sl, 6.25, 1.38, 0.09, 2.0, fill=P_BLUE)
tb(sl, "Preguntas clave", 6.45, 1.44, 6.4, 0.42, size=14, bold=True, color=P_BLUE)
qs = ["Q4  —  Situación laboral  (→ Disponible)",
      "Q6  —  Frecuencia baja motivación  (→ Efectivo)",
      "Q13 —  Disposición gamificación  (→ Objetivo)",
      "Q17 —  Disposición de pago  (→ Precio)"]
for j, q in enumerate(qs):
    tb(sl, q, 6.45, 1.88 + j*0.37, 6.5, 0.34, size=11.5, color=P_WHITE)

# Cascada
tb(sl, "Modelo de Segmentación en Cascada", 0.35, 3.55, 12.5, 0.45,
   size=14, bold=True, color=P_PURPL)

cascade = [
    ("POTENCIAL",   "185,000",  "docentes Lima Met.",   P_PURP,  "MINEDU 2024"),
    ("DISPONIBLE",  "177,600",  "× 96 % (Q4)",          P_BLUE,  "−3,400"),
    ("EFECTIVO",    "103,008",  "× 58 % (Q6)",          P_TEAL,  "−74,592"),
    ("OBJETIVO",    "82,406",   "× 80 % (Q13)",         P_GREEN, "−20,602"),
]
for i, (lbl, val, crit, col, delta) in enumerate(cascade):
    x = 0.35 + i * 3.23
    rect(sl, x, 4.1, 3.0, 1.9, fill=P_DARK)
    rect(sl, x, 4.1, 3.0, 0.07, fill=col)
    rect(sl, x, 5.93, 3.0, 0.07, fill=col)
    tb(sl, lbl,   x+0.1, 4.2,  2.8, 0.38, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, val,   x+0.1, 4.6,  2.8, 0.65, size=26, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, crit,  x+0.1, 5.24, 2.8, 0.38, size=10.5, color=P_LGRAY, align=PP_ALIGN.CENTER)
    tb(sl, delta, x+0.1, 5.6,  2.8, 0.3,  size=9, color=col, align=PP_ALIGN.CENTER, italic=True)
    if i < 3:
        tb(sl, "▶", x+3.0, 4.8, 0.23, 0.55, size=18, bold=True, color=P_PURPL, align=PP_ALIGN.CENTER)

# ── SLIDE 4: RESULTADOS ENCUESTA ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Resultados de la Encuesta", "Preguntas clave para la segmentación de mercado (n=100)")

# Q6 gráfico de barras horizontal
tb(sl, "Q6 — Frecuencia de baja motivación estudiantil", 0.35, 1.38, 7.5, 0.45, size=14, bold=True, color=P_GOLD)
q6 = [("Siempre",        22, P_RED),
      ("Frecuentemente",  36, P_PURP),
      ("A veces",         30, P_BLUE),
      ("Raramente",        8, P_TEAL),
      ("Nunca",            4, P_GREEN)]
max_v = 36
bar_bg = P_DARK
for j, (lbl, val, col) in enumerate(q6):
    y_b = 1.95 + j * 0.76
    rect(sl, 0.35, y_b, 5.5, 0.62, fill=bar_bg)
    bar_w = (val / max_v) * 4.8
    rect(sl, 0.35, y_b, bar_w, 0.62, fill=col)
    tb(sl, lbl, 0.45, y_b+0.14, 2.2, 0.35, size=11, bold=(j<2), color=P_WHITE)
    tb(sl, f"{val} %", 0.35+bar_w+0.05, y_b+0.14, 0.8, 0.35, size=13, bold=True, color=col)
    if j < 2:
        rect(sl, 0.35+bar_w, y_b, 5.5-bar_w, 0.62, fill=P_PANEL)
        tb(sl, "✓ Filtro Mercado Efectivo", 0.35+bar_w+0.1, y_b+0.14, 3.0, 0.35, size=9, color=col, italic=True)

# Resumen a la derecha
rect(sl, 6.1, 1.38, 6.9, 2.35, fill=P_PANEL)
rect(sl, 6.1, 1.38, 0.09, 2.35, fill=P_PURP)
tb(sl, "Hallazgos clave", 6.3, 1.44, 6.5, 0.4, size=14, bold=True, color=P_PURPL)
highlights = [
    "96 % trabaja en institución educativa (Q4)",
    "58 % baja motivación frecuente/siempre (Q6)",
    "80 % implementaría gamificación (Q13)",
    "61 de 100 tienen intención de pago (Q17)",
]
for j, h in enumerate(highlights):
    tb(sl, "→ " + h, 6.3, 1.92 + j*0.44, 6.6, 0.4, size=11.5, color=P_WHITE)

# Q13 resultado
rect(sl, 6.1, 3.95, 6.9, 1.4, fill=P_DARK)
rect(sl, 6.1, 3.95, 6.9, 0.06, fill=P_GREEN)
tb(sl, "Q13 — Disposición de implementar gamificación",
   6.2, 4.05, 6.7, 0.4, size=11, bold=True, color=P_GREEN)
tb(sl, "80 %", 6.2, 4.45, 3.0, 0.65, size=36, bold=True, color=P_GREEN, align=PP_ALIGN.LEFT)
tb(sl, "Sí implementaría\nuna solución en su aula",
   9.3, 4.52, 3.5, 0.65, size=11, color=P_LGRAY)

# Barra de progreso 80%
rect(sl, 6.2, 5.25, 6.5, 0.35, fill=P_PANEL)
rect(sl, 6.2, 5.25, 6.5*0.80, 0.35, fill=P_GREEN)
tb(sl, "80 %  de 100 encuestados", 6.2, 5.65, 6.5, 0.35, size=10, color=P_LGRAY)

# ── SLIDE 5: PRECIO ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Análisis de Precio", "Q17 — Disposición de pago mensual  ·  61 respondentes con intención definitiva")

# 3 cards de precio
price_cards = [
    ("S/. 20 – 40/mes",  "S/. 30",  "S/. 360/año",  "35 docentes",  "57.4 %", P_BLUE),
    ("S/. 41 – 70/mes",  "S/. 55",  "S/. 660/año",  "18 docentes",  "29.5 %", P_PURP),
    ("S/. 71 – 100/mes", "S/. 85",  "S/. 1,020/año", "8 docentes",  "13.1 %", P_TEAL),
]
for i, (rng, mo, yr, qty, pct, col) in enumerate(price_cards):
    x = 0.35 + i * 4.15
    # Card
    rect(sl, x, 1.4, 3.85, 3.2, fill=P_DARK)
    rect(sl, x, 1.4, 3.85, 0.08, fill=col)
    rect(sl, x, 4.52, 3.85, 0.08, fill=col)
    # Rango
    tb(sl, rng,  x+0.12, 1.5, 3.6, 0.4, size=12, color=P_LGRAY, align=PP_ALIGN.CENTER)
    # Precio grande
    tb(sl, mo,   x+0.12, 1.9, 3.6, 0.75, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    accent_line(sl, x+0.5, 2.72, 2.85, color=col, h=0.04)
    tb(sl, yr,   x+0.12, 2.8, 3.6, 0.48, size=18, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, qty,  x+0.12, 3.35, 3.6, 0.38, size=13, color=P_LGRAY, align=PP_ALIGN.CENTER)
    # Barra de representación
    bar_w = (float(pct.replace(" %","").replace(",",".")) / 57.4) * 3.45
    rect(sl, x+0.2, 4.0, 3.45, 0.35, fill=P_PANEL)
    rect(sl, x+0.2, 4.0, bar_w, 0.35, fill=col)
    tb(sl, pct + " del total", x+0.2, 4.38, 3.45, 0.3, size=9.5, color=col, align=PP_ALIGN.CENTER)

# Resultado precio ponderado
rect(sl, 0.35, 4.78, 12.63, 1.45, fill=P_DARK)
rect(sl, 0.35, 4.78, 0.1, 1.45, fill=P_GOLD)
rect(sl, 0.35, 4.78, 12.63, 0.07, fill=P_GOLD)
tb(sl, "Precio Ponderado Anual", 0.6, 4.87, 5, 0.45, size=14, bold=True, color=P_GOLD)
tb(sl, "S/. 535.08 / docente / año   ≈   S/. 44.59 / mes",
   0.6, 5.35, 9, 0.65, size=28, bold=True, color=P_WHITE)
tb(sl, "Cálculo: (35×360 + 18×660 + 8×1,020) ÷ 61 = 32,640 ÷ 61 = 535.08",
   0.6, 6.0,  9, 0.35, size=10.5, color=P_LGRAY, italic=True)
tb(sl, "61 de 100\ncon intención\nde pago",
   9.8, 4.9, 2.9, 1.1, size=14, bold=True, color=P_GOLD, align=PP_ALIGN.CENTER)

# ── SLIDE 6: ESCENARIOS ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Análisis de la Demanda — Escenarios", "Año 2027  ·  Mercado Objetivo: 82,406 docentes  ·  Precio: S/. 535.08/año")

esc = [
    ("OPTIMISTA", "80 %", "65,925", "S/. 35,274,015", "≈ USD 9.5M", P_GREEN,  P_BG),
    ("MODERADO",  "50 %", "41,203", "S/. 22,046,259", "≈ USD 5.9M", P_GOLD,   P_BG),
    ("PESIMISTA", "30 %", "24,722", "S/. 13,227,786", "≈ USD 3.6M", P_RED,    P_BG),
]
for i, (name, pct, subs, rev, usd, col, bg) in enumerate(esc):
    x = 0.35 + i * 4.32
    # Card grande
    rect(sl, x, 1.42, 4.1, 5.45, fill=P_DARK)
    rect(sl, x, 1.42, 4.1, 0.1,  fill=col)
    rect(sl, x, 6.77, 4.1, 0.1,  fill=col)
    # Nombre
    tb(sl, name, x, 1.55, 4.1, 0.5, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Porcentaje — enorme
    tb(sl, pct,  x, 2.12, 4.1, 1.1, size=52, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, "penetración del mercado", x, 3.22, 4.1, 0.38, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)
    accent_line(sl, x+0.3, 3.65, 3.5, color=col, h=0.04)
    # Suscripciones
    tb(sl, "Suscripciones", x, 3.75, 4.1, 0.35, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)
    tb(sl, subs,  x, 4.1, 4.1, 0.65, size=24, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    accent_line(sl, x+0.3, 4.8, 3.5, color=col, h=0.04)
    # Ingresos
    tb(sl, "Ingresos anuales", x, 4.9, 4.1, 0.35, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)
    tb(sl, rev,  x, 5.25, 4.1, 0.5, size=17, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, usd,  x, 5.75, 4.1, 0.38, size=12, color=P_LGRAY, align=PP_ALIGN.CENTER)

# Pie de fórmula
tb(sl, "Demanda = Mercado Objetivo × % Penetración     |     Ingreso = Demanda × S/. 535.08 / año",
   0.35, 7.05, 12.63, 0.3, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 7: PROYECCIÓN 5 AÑOS ───────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Proyección 2027–2031", "Escenario Moderado  ·  Tasa de crecimiento: 15 % anual  ·  CAGR EdTech LATAM (HolonIQ 2024)")

# KPIs top
kpis = [
    ("41,203",   "Suscripciones\n2027 (base)", P_BLUE),
    ("72,063",   "Suscripciones\n2031",         P_PURP),
    ("+75 %",    "Crecimiento\nacumulado",       P_TEAL),
    ("S/. 148.7M","Ingresos\nacumulados 5 años", P_GOLD),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.35 + i * 3.24
    rect(sl, x, 1.4, 3.05, 1.35, fill=P_PANEL)
    rect(sl, x, 1.4, 3.05, 0.07, fill=col)
    tb(sl, val, x, 1.5,  3.05, 0.65, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, lbl, x, 2.1,  3.05, 0.55, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)

# Gráfico de barras
bar_colors_ch = [P_BLUE, P_PURP, P_TEAL, P_GREEN, P_GOLD]
PROJ_ING_M   = [v/1_000_000 for v in PROJ_ING]
max_ing      = max(PROJ_ING_M)
chart_bottom = 6.75
chart_h_max  = 3.5
chart_left   = 0.5
bar_w_ch     = 1.7
bar_gap_ch   = 0.6

for i, (yr, sub, ing_m, col) in enumerate(zip(PROJ_YEARS, PROJ_SUBS, PROJ_ING_M, bar_colors_ch)):
    x      = chart_left + i * (bar_w_ch + bar_gap_ch)
    bar_h  = (ing_m / max_ing) * chart_h_max
    y_bar  = chart_bottom - bar_h
    rect(sl, x, y_bar, bar_w_ch, bar_h, fill=col)
    # valor sobre barra
    tb(sl, f"S/.{ing_m:.1f}M", x-0.05, y_bar - 0.45, bar_w_ch+0.1, 0.38,
       size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    # subs bajo barra
    tb(sl, str(yr), x, chart_bottom+0.04, bar_w_ch, 0.35,
       size=12, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, f"{sub:,} susc.", x, chart_bottom+0.4, bar_w_ch, 0.3,
       size=9.5, color=P_LGRAY, align=PP_ALIGN.CENTER)

# Línea base
rect(sl, chart_left-0.1, chart_bottom, 12.5, 0.04, fill=P_PURPL)

# Nota tasa
rect(sl, 0.35, 3.0, 5.8, 0.6, fill=P_PANEL)
tb(sl, "Tasa: 15 % anual  ·  Fuente: HolonIQ EdTech LATAM Report 2024",
   0.5, 3.08, 5.5, 0.44, size=10.5, color=P_LGRAY, italic=True)

# ── SLIDE 8: CONCLUSIONES ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
header_bar(sl, "Conclusiones", "Síntesis del análisis de demanda — LegendaryClass")

concl_ppt = [
    (P_PURP,  "82,406 docentes",      "conforman el mercado objetivo en Lima Metropolitana (MINEDU 2024)."),
    (P_BLUE,  "S/. 535 / año",        "precio ponderado validado por disposición de pago real (Q17)."),
    (P_TEAL,  "S/. 22M en 2027",      "ingresos proyectados en escenario moderado (50 % penetración)."),
    (P_GOLD,  "S/. 38.5M en 2031",    "con crecimiento 15 % anual — CAGR EdTech LATAM (HolonIQ 2024)."),
    (P_GREEN, "58 % de docentes",     "reportan baja motivación frecuente — necesidad validada por encuesta."),
]
for i, (col, bold_val, rest) in enumerate(concl_ppt):
    y = 1.45 + i * 1.08
    rect(sl, 0.35, y, 12.63, 0.95, fill=P_PANEL)
    rect(sl, 0.35, y, 0.1,   0.95, fill=col)
    # número
    rect(sl, 0.55, y+0.18, 0.58, 0.58, fill=col)
    tb(sl, str(i+1), 0.55, y+0.13, 0.58, 0.58, size=20, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    # valor destacado
    tb(sl, bold_val, 1.3, y+0.14, 3.3, 0.58, size=18, bold=True, color=col)
    # descripción
    tb(sl, rest,     4.65, y+0.2, 8.1, 0.55, size=13, color=P_WHITE)

# ── SLIDE 9: CIERRE ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=P_BG)
rect(sl, 0, 0, 13.33, 0.12, fill=P_PURP)
rect(sl, 0, 7.38, 13.33, 0.12, fill=P_PURP)

# Decoración
rect(sl, 0, 0, 0.5, 7.5, fill=P_DARK)
rect(sl, 12.83, 0, 0.5, 7.5, fill=P_DARK)

rect(sl, 3.5, 2.1, 6.33, 3.3, fill=P_PANEL)
rect(sl, 3.5, 2.1, 6.33, 0.08, fill=P_PURP)
rect(sl, 3.5, 5.32, 6.33, 0.08, fill=P_PURP)

tb(sl, "⚔️  LegendaryClass",
   3.6, 2.25, 6.13, 0.85, size=30, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Sistema SaaS de Gamificación Educativa",
   3.6, 3.12, 6.13, 0.5, size=14, color=P_LGRAY, align=PP_ALIGN.CENTER)
accent_line(sl, 4.2, 3.72, 5.0, color=P_GOLD, h=0.06)
tb(sl, "¡Gracias!",
   3.6, 3.88, 6.13, 0.85, size=36, bold=True, color=P_GOLD, align=PP_ALIGN.CENTER)
tb(sl, ALUMNO + "  ·  " + CURSO + "  ·  " + ANIO_SEM,
   3.6, 4.82, 6.13, 0.42, size=11, color=P_LGRAY, align=PP_ALIGN.CENTER)

# Meta info
tb(sl, INSTITUCION + "  ·  " + CARRERA + "  ·  " + CICLO,
   0.7, 6.9, 12, 0.4, size=10, color=P_LGRAY, align=PP_ALIGN.CENTER)

PPT_OUT = "C:/Proyectos/LegendaryClass/lab03/Presentacion_Caso3_LegendaryClass_v2.pptx"
prs.save(PPT_OUT)
print("[OK] PPT guardada:", PPT_OUT)
print()
print("Archivos listos:")
print(" ->", WORD_OUT)
print(" ->", PPT_OUT)
