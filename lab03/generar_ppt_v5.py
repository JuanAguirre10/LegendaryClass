"""
PPT v5 - Estilo fiel al PHP original:
- Fondo: fondo.png + overlay blanco (igual que la web)
- Título: gradiente dorado REAL via XML (igual que .super-clear-title / .epic-title)
- Sin cards oscuras — todo sobre fondo claro
- Borde/acento: dorado #fbbf24
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta del PHP ───────────────────────────────────────
GOLD_400  = RGB(251,191, 36)   # #fbbf24
GOLD_500  = RGB(245,158, 11)   # #f59e0b
GOLD_600  = RGB(217,119,  6)   # #d97706
GOLD_800  = RGB(146, 64, 14)   # #92400e
NAVY      = RGB( 15, 23, 42)   # #0f172a
PURPLE_D  = RGB( 88, 28,135)   # #581c87
PURPLE_M  = RGB(124, 58,237)   # #7c3aed
PURPLE_L  = RGB(196,181,253)   # purple-300
GRAY_800  = RGB( 31, 41, 55)
GRAY_600  = RGB( 75, 85, 99)
GRAY_400  = RGB(156,163,175)
GRAY_200  = RGB(229,231,235)
WHITE     = RGB(255,255,255)
CREAM     = RGB(255,253,235)

FONDO = "C:/xampp/htdocs/LegendaryClass/public/fondo.png"

INSTITUCION = "TECSUP"
CARRERA     = "Ingeniería de Software  ·  V Ciclo"
CURSO       = "Diseño de Proyectos de Innovación"
DOCENTE     = "Turkowsky Vizcarra, Luisa"
ALUMNO      = "Aguirre, Juan"
ANIO        = "2026 - I"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════

def _rect(sl, l, t, w, h, fill=None, line_color=None, line_w=9525):
    sp = sl.shapes.add_shape(1,
        Inches(max(l,-1)), Inches(max(t,-1)),
        Inches(min(w,15)),  Inches(min(h,11)))
    sp.fill.background() if fill is None else (
        sp.fill.solid().__class__, sp.fill.solid(), setattr(sp.fill.fore_color,'rgb',fill))
    sp.line.fill.background() if line_color is None else (
        setattr(sp.line.color,'rgb',line_color), setattr(sp.line,'width',Emu(line_w)))
    return sp

def rect(sl, l, t, w, h, fill=None, line_color=None, line_w=9525):
    return _rect(sl, l, t, w, h, fill, line_color, line_w)

def txt(sl, text, l, t, w, h,
        size=13, bold=False, color=GRAY_800,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def gradient_title(sl, text, l, t, w, h, size=48):
    """
    Título con gradiente dorado igual al CSS:
    linear-gradient(45deg, #fbbf24, #f59e0b, #d97706, #92400e, #fbbf24)
    Implementado via XML directo en el run de texto.
    """
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER

    # Espaciado de letras (letter-spacing: 0.15em aprox)
    pPr = p._p.get_or_add_pPr()

    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True

    # Acceder al XML del run para aplicar gradiente de relleno
    rPr = run._r.get_or_add_rPr()

    # Eliminar cualquier solidFill existente
    for el in rPr.findall(qn('a:solidFill')):
        rPr.remove(el)

    # Construir gradFill
    gradFill = etree.SubElement(rPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    stops = [
        (0,      'fbbf24'),
        (25000,  'f59e0b'),
        (50000,  'd97706'),
        (75000,  '92400e'),
        (100000, 'fbbf24'),
    ]
    for pos, hex_color in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', hex_color)

    # Ángulo 45° = 2700000 en unidades OOXML (60000 * grados)
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', '2700000')
    lin.set('scaled', '0')

    return box

def add_bg_image(sl, opacity_pct=82):
    """Fondo con fondo.png + overlay blanco semitransparente"""
    if os.path.exists(FONDO):
        sl.shapes.add_picture(FONDO, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    # Overlay blanco
    sp = rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    spPr = sp._element.find(qn('p:spPr'))
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None:
        sg = sf.find(qn('a:srgbClr'))
        if sg is None: sg = etree.SubElement(sf, qn('a:srgbClr'))
        sg.set('val', 'FFFFFF')
        al = etree.SubElement(sg, qn('a:alpha'))
        al.set('val', str(int(opacity_pct * 1000)))

def gold_line(sl, l, t, w, h=0.05):
    rect(sl, l, t, w, h, fill=GOLD_500)

def card_white(sl, l, t, w, h):
    """Card blanca con borde dorado suave"""
    rect(sl, l, t, w, h, fill=WHITE, line_color=GOLD_500, line_w=15000)

def footer(sl):
    gold_line(sl, 0, 7.3, 13.33, h=0.03)
    txt(sl, f"LegendaryClass  ·  {CURSO}  ·  {INSTITUCION}  ·  {ANIO}",
        0.4, 7.36, 12.5, 0.2, size=8, color=GRAY_400)

def slide_header(sl, title, subtitle=""):
    """
    Header igual al epic-header del PHP pero SIN fondo oscuro —
    solo borde dorado inferior y título en dorado.
    """
    gold_line(sl, 0, 1.22, 13.33, h=0.06)
    gradient_title(sl, title.upper(), 0.4, 0.1, 12.5, 0.72, size=22)
    if subtitle:
        txt(sl, subtitle, 0.4, 0.82, 12.5, 0.38,
            size=11, color=GRAY_600, italic=True)
    footer(sl)


# ════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, opacity_pct=80)

# Título principal — igual al .super-clear-title del PHP
gradient_title(sl, "LEGENDARYCLASS", 0.5, 0.5, 12.33, 1.55, size=62)

txt(sl, "⚔️  Bienvenido al Portal de Aventureros  ⚔️",
    0.5, 2.1, 12.33, 0.5, size=14, color=GRAY_600,
    align=PP_ALIGN.CENTER, italic=True)

gold_line(sl, 3.0, 2.72, 7.33)

# Panel central (estilo .welcome-panel del PHP)
card_white(sl, 3.0, 2.85, 7.33, 3.75)

txt(sl, CURSO, 3.2, 3.0, 6.93, 0.45,
    size=13, bold=True, color=GRAY_800, align=PP_ALIGN.CENTER)
txt(sl, "Caso 3 — Parte Práctica  ·  Determinación de la Demanda",
    3.2, 3.48, 6.93, 0.42, size=12, color=GOLD_600, align=PP_ALIGN.CENTER)

gold_line(sl, 3.5, 4.02, 6.33)

txt(sl, f"Docente:   {DOCENTE}",
    3.2, 4.12, 6.93, 0.38, size=11, color=GRAY_600, align=PP_ALIGN.CENTER)
txt(sl, f"Alumno:    {ALUMNO}",
    3.2, 4.5,  6.93, 0.38, size=11, color=GRAY_600, align=PP_ALIGN.CENTER)

gold_line(sl, 3.5, 4.98, 6.33)

txt(sl, f"{INSTITUCION}  ·  {CARRERA}  ·  {ANIO}",
    3.2, 5.08, 6.93, 0.38, size=10.5, color=GRAY_400, align=PP_ALIGN.CENTER)

footer(sl)


# ════════════════════════════════════════════════════════
# SLIDE 2 — DESCRIPCIÓN DEL PRODUCTO
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Descripción del Producto",
             "LegendaryClass — plataforma SaaS de gamificación educativa para docentes")

txt(sl, "Transforma el aula en un RPG educativo. Los estudiantes eligen clase de personaje, "
        "acumulan XP, suben de nivel y canjean recompensas. El docente gestiona comportamientos, "
        "misiones y reportes de progreso.",
    0.45, 1.38, 12.43, 0.72, size=12.5, color=GRAY_600, wrap=True)

features = [
    ("⚔️", "5 Clases RPG",     "Mago · Guerrero · Ninja\nArquero · Lanzador\n+20 % XP por acción"),
    ("📈", "Motor de XP",       "Niveles, racha diaria\ny logros automáticos"),
    ("👩‍🏫","Panel Docente",    "Comportamientos, recompensas\ny reportes por aula"),
    ("🏛️","Panel Director",   "Estadísticas globales\nde la institución"),
]
for i, (ico, tit, desc) in enumerate(features):
    x = 0.45 + i * 3.23
    card_white(sl, x, 2.25, 3.05, 2.95)
    gold_line(sl, x, 2.25, 3.05)
    txt(sl, ico,  x+0.15, 2.35, 0.65, 0.55, size=26)
    txt(sl, tit,  x+0.15, 2.95, 2.75, 0.45, size=12.5, bold=True, color=GOLD_600)
    txt(sl, desc, x+0.15, 3.42, 2.75, 1.6,  size=10.5, color=GRAY_600, wrap=True)

gold_line(sl, 0.45, 5.38, 12.43)
txt(sl, "Stack:   Angular 18  ·  NestJS  ·  PostgreSQL  ·  Prisma  ·  TailwindCSS",
    0.45, 5.48, 9, 0.38, size=11.5, bold=True, color=GRAY_800)
txt(sl, "Precio:  S/. 30 – 85 / mes por docente   ·   Sin costo para estudiantes",
    0.45, 5.88, 10, 0.35, size=11, color=GRAY_600)


# ════════════════════════════════════════════════════════
# SLIDE 3 — METODOLOGÍA Y SEGMENTACIÓN
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Metodología", "Encuesta a 100 docentes de Lima Metropolitana  ·  feb–mar 2026")

card_white(sl, 0.45, 1.42, 5.7, 1.68)
gold_line(sl, 0.45, 1.42, 5.7)
txt(sl, "Encuesta de campo", 0.65, 1.52, 5.3, 0.42, size=12.5, bold=True, color=GOLD_600)
for j, item in enumerate([
    "100 docentes · Lima Metropolitana",
    "Google Forms · 20 preguntas",
    "Periodo: febrero – marzo 2026",
]):
    txt(sl, "·  " + item, 0.65, 1.98+j*0.36, 5.3, 0.33, size=11.5, color=GRAY_600)

card_white(sl, 6.45, 1.42, 6.43, 1.68)
gold_line(sl, 6.45, 1.42, 6.43)
txt(sl, "Preguntas clave", 6.65, 1.52, 6.1, 0.42, size=12.5, bold=True, color=GOLD_600)
for j, q in enumerate([
    "Q4  — Situación laboral  →  Disponible",
    "Q6  — Baja motivación  →  Efectivo",
    "Q13 — Disposición gamificación  →  Objetivo",
    "Q17 — Disposición de pago  →  Precio",
]):
    txt(sl, q, 6.65, 1.98+j*0.36, 6.1, 0.33, size=11, color=GRAY_600)

txt(sl, "Segmentación en Cascada", 0.45, 3.28, 8, 0.42,
    size=13, bold=True, color=GRAY_800)

cascade = [
    ("POTENCIAL",  "185,000", "MINEDU 2024",  GRAY_200, GRAY_800, GOLD_600),
    ("DISPONIBLE", "177,600", "× 96 %  (Q4)", CREAM,    GRAY_800, GOLD_600),
    ("EFECTIVO",   "103,008", "× 58 %  (Q6)", CREAM,    GRAY_800, GOLD_600),
    ("OBJETIVO",    "82,406", "× 80 %  (Q13)",GOLD_500, WHITE,    WHITE),
]
for i, (lbl, val, crit, bg, fg, sub_fg) in enumerate(cascade):
    x = 0.45 + i * 3.23
    rect(sl, x, 3.82, 3.05, 2.35, fill=bg, line_color=GOLD_500, line_w=12700)
    gold_line(sl, x, 3.82, 3.05)
    txt(sl, lbl,  x+0.15, 3.92, 2.75, 0.38, size=10, bold=True, color=fg)
    txt(sl, val,  x+0.15, 4.33, 2.75, 0.65, size=26, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(sl, crit, x+0.15, 5.0,  2.75, 0.35, size=10.5, color=sub_fg, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "›", x+3.05, 4.55, 0.18, 0.5, size=20, bold=True,
            color=GOLD_500, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS DE ENCUESTA
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Resultados de la Encuesta",
             "Preguntas clave para la segmentación  ·  n = 100")

txt(sl, "Q6  —  Frecuencia de baja motivación estudiantil",
    0.45, 1.42, 7.5, 0.42, size=13, bold=True, color=GRAY_800)

q6 = [("Siempre",        22, GOLD_600, True),
      ("Frecuentemente",  36, GOLD_500, True),
      ("A veces",         30, GRAY_200, False),
      ("Raramente",        8, GRAY_200, False),
      ("Nunca",            4, GRAY_200, False)]
for j, (lbl, val, col, hi) in enumerate(q6):
    y = 1.95 + j * 0.72
    rect(sl, 0.45, y, 7.0, 0.58, fill=GRAY_200)
    rect(sl, 0.45, y, (val/36)*6.5, 0.58, fill=col)
    txt(sl, lbl, 0.58, y+0.13, 2.2, 0.3,
        size=11, bold=hi, color=WHITE if hi else GRAY_600)
    txt(sl, f"{val} %", 0.45+(val/36)*6.5+0.1, y+0.13, 0.7, 0.3,
        size=12, bold=hi, color=GOLD_600 if hi else GRAY_400)
    if hi:
        txt(sl, "← filtro", 0.45+(val/36)*6.5+0.85, y+0.15, 1.5, 0.28,
            size=9, color=GOLD_600, italic=True)

card_white(sl, 0.45, 5.62, 7.0, 0.68)
gold_line(sl, 0.45, 5.62, 7.0)
txt(sl, "Siempre + Frecuentemente  =  58 %   →   Mercado Efectivo",
    0.65, 5.74, 6.5, 0.42, size=12, bold=True, color=GOLD_600)

txt(sl, "Resumen por pregunta", 7.8, 1.42, 5.1, 0.42, size=13, bold=True, color=GRAY_800)
hallazgos = [
    ("Q4",  "96 %",  "Trabajan en institución educativa",   GOLD_600),
    ("Q6",  "58 %",  "Baja motivación frecuente/siempre",   GOLD_500),
    ("Q13", "80 %",  "Implementarían gamificación",          GRAY_800),
    ("Q17", "61/100","Con intención definitiva de pago",     GRAY_600),
]
for j, (q, val, desc, col) in enumerate(hallazgos):
    y = 1.95 + j * 1.12
    card_white(sl, 7.8, y, 5.1, 0.95)
    gold_line(sl, 7.8, y, 5.1)
    txt(sl, q,   7.95, y+0.12, 0.6, 0.38, size=10, bold=True, color=GOLD_500)
    txt(sl, val, 8.6,  y+0.08, 1.8, 0.5,  size=22, bold=True, color=col)
    txt(sl, desc, 10.4, y+0.18, 2.3, 0.5, size=10, color=GRAY_600, wrap=True)


# ════════════════════════════════════════════════════════
# SLIDE 5 — ANÁLISIS DE PRECIO
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Análisis de Precio",
             "Q17 — Disposición de pago  ·  61 respondentes con intención definitiva")

price_cards = [
    ("S/. 20 – 40 / mes",  "S/. 30",  "S/. 360 / año",  35, GRAY_800),
    ("S/. 41 – 70 / mes",  "S/. 55",  "S/. 660 / año",  18, GOLD_600),
    ("S/. 71 – 100 / mes", "S/. 85", "S/. 1,020 / año",  8, GOLD_800),
]
for i, (rng, mo, yr, qty, col) in enumerate(price_cards):
    x = 0.45 + i * 4.15
    card_white(sl, x, 1.42, 3.9, 3.7)
    gold_line(sl, x, 1.42, 3.9)
    txt(sl, rng, x+0.15, 1.52, 3.6, 0.38, size=10.5, color=GRAY_400, align=PP_ALIGN.CENTER)
    txt(sl, mo,  x+0.15, 1.95, 3.6, 0.85, size=38, bold=True, color=col, align=PP_ALIGN.CENTER)
    gold_line(sl, x+0.3, 2.88, 3.3)
    txt(sl, yr,  x+0.15, 2.98, 3.6, 0.5,  size=18, bold=True, color=GOLD_600, align=PP_ALIGN.CENTER)
    txt(sl, f"{qty} docentes", x+0.15, 3.55, 3.6, 0.38,
        size=12, color=GRAY_400, align=PP_ALIGN.CENTER)

# Resultado — panel dorado
rect(sl, 0.45, 5.3, 12.43, 1.28, fill=CREAM, line_color=GOLD_500, line_w=18000)
gold_line(sl, 0.45, 5.3, 12.43)
txt(sl, "Precio Ponderado Anual", 0.7, 5.38, 5.5, 0.42, size=11.5, bold=True, color=GOLD_600)
txt(sl, "S/. 535.08 / docente / año   ≈   S/. 44.59 / mes",
    0.7, 5.8, 11, 0.55, size=26, bold=True, color=GRAY_800)
txt(sl, "( 35 × S/.360  +  18 × S/.660  +  8 × S/.1,020 )  ÷  61  =  S/.32,640 ÷ 61",
    0.7, 6.38, 11.5, 0.28, size=10, color=GRAY_400, italic=True)


# ════════════════════════════════════════════════════════
# SLIDE 6 — ESCENARIOS
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Escenarios de Demanda",
             "Año 2027  ·  Mercado Objetivo: 82,406 docentes  ·  Precio: S/. 535.08 / año")

esc = [
    ("Optimista", "80 %", "65,925", "S/. 35,274,015", CREAM,    GRAY_800, GOLD_600),
    ("Moderado",  "50 %", "41,203", "S/. 22,046,259", GRAY_200, GRAY_800, GOLD_600),
    ("Pesimista", "30 %", "24,722", "S/. 13,227,786", WHITE,    GRAY_600, GRAY_600),
]
for i, (name, pct, subs, rev, bg, fg, acc) in enumerate(esc):
    x = 0.45 + i * 4.3
    rect(sl, x, 1.42, 4.08, 5.45, fill=bg, line_color=GOLD_500, line_w=15000)
    gold_line(sl, x, 1.42, 4.08)
    txt(sl, name,  x+0.2, 1.52, 3.7, 0.45, size=15, bold=True, color=acc)
    txt(sl, pct,   x+0.2, 2.0,  3.7, 1.1,  size=54, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(sl, "penetración", x+0.2, 3.15, 3.7, 0.35, size=10, color=GRAY_400, align=PP_ALIGN.CENTER)
    gold_line(sl, x+0.3, 3.6, 3.48)
    txt(sl, "Suscripciones", x+0.2, 3.72, 3.7, 0.35, size=10, color=GRAY_400, align=PP_ALIGN.CENTER)
    txt(sl, subs,  x+0.2, 4.08, 3.7, 0.62, size=24, bold=True, color=fg, align=PP_ALIGN.CENTER)
    gold_line(sl, x+0.3, 4.78, 3.48)
    txt(sl, "Ingresos anuales", x+0.2, 4.9, 3.7, 0.35, size=10, color=GRAY_400, align=PP_ALIGN.CENTER)
    txt(sl, rev,   x+0.2, 5.28, 3.7, 0.52, size=14.5, bold=True, color=acc, align=PP_ALIGN.CENTER)

txt(sl, "Demanda = Mercado Objetivo × % Penetración   ·   Ingreso = Demanda × S/. 535.08 / año",
    0.45, 7.05, 12.43, 0.28, size=9.5, color=GRAY_400, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 7 — PROYECCIÓN 5 AÑOS
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Proyección 2027–2031",
             "Escenario Moderado  ·  Crecimiento anual: +15 %  ·  HolonIQ EdTech LATAM 2024")

kpis = [
    ("41,203 → 72,063", "Suscripciones  2027 → 2031", GOLD_600),
    ("+75 %",           "Crecimiento acumulado",        GRAY_800),
    ("S/. 148.7M",      "Ingresos acumulados 5 años",   GOLD_600),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.45 + i * 4.3
    rect(sl, x, 1.42, 4.08, 1.2, fill=CREAM, line_color=GOLD_500, line_w=12700)
    gold_line(sl, x, 1.42, 4.08)
    txt(sl, val, x+0.2, 1.5,  3.7, 0.6, size=19, bold=True, color=col)
    txt(sl, lbl, x+0.2, 2.1,  3.7, 0.4, size=10, color=GRAY_600)

PROJ_YEARS = [2027, 2028, 2029, 2030, 2031]
PROJ_SUBS  = [41203, 47383, 54490, 62664, 72063]
PROJ_ING   = [22046259, 25353198, 29156177, 33529604, 38559045]
bar_cols   = [GRAY_200, CREAM, GOLD_500, GOLD_500, GOLD_400]
bar_line   = [GOLD_600, GOLD_600, GOLD_600, GOLD_600, GOLD_600]

max_ing   = max(PROJ_ING)
chart_b   = 6.6
chart_h   = 3.25
bw        = 1.72; gap = 0.48; sx = 0.8

gold_line(sl, sx-0.1, chart_b, 12.1, h=0.04)

for i, (yr, sub, ing, col, lc) in enumerate(zip(PROJ_YEARS, PROJ_SUBS, PROJ_ING, bar_cols, bar_line)):
    x  = sx + i*(bw+gap)
    bh = (ing/max_ing)*chart_h
    yb = chart_b - bh
    rect(sl, x, yb, bw, bh, fill=col, line_color=lc, line_w=9525)
    ing_m = ing/1_000_000
    txt(sl, f"S/.{ing_m:.1f}M", x-0.05, yb-0.42, bw+0.1, 0.35,
        size=11, bold=True, color=GOLD_600, align=PP_ALIGN.CENTER)
    txt(sl, str(yr), x, chart_b+0.07, bw, 0.33,
        size=12, bold=True, color=GRAY_800, align=PP_ALIGN.CENTER)
    txt(sl, f"{sub:,}", x, chart_b+0.42, bw, 0.28,
        size=9, color=GRAY_400, align=PP_ALIGN.CENTER)

txt(sl, "HolonIQ EdTech LATAM Report 2024  ·  CAGR 15 % para plataformas SaaS educativas en Perú y LATAM",
    0.45, 7.1, 12.43, 0.25, size=8, color=GRAY_400, italic=True)


# ════════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSIONES
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 88)
slide_header(sl, "Conclusiones", "Síntesis del análisis de demanda — LegendaryClass")

concl = [
    ("82,406 docentes",   "conforman el mercado objetivo en Lima Metropolitana (MINEDU 2024).",  GOLD_600),
    ("S/. 535 / año",     "precio ponderado validado por disposición de pago real (Q17).",        GOLD_600),
    ("S/. 22M en 2027",   "ingresos proyectados en escenario moderado (50 % penetración).",       GRAY_800),
    ("S/. 38.5M en 2031", "con crecimiento anual del 15 % — CAGR EdTech LATAM (HolonIQ 2024).", GRAY_800),
    ("58 % de docentes",  "reportan baja motivación frecuente — necesidad real y cuantificada.",  GOLD_600),
]
for i, (val, rest, col) in enumerate(concl):
    y = 1.42 + i * 1.03
    card_white(sl, 0.45, y, 12.43, 0.9)
    gold_line(sl, 0.45, y, 12.43)
    rect(sl, 0.45, y, 0.1, 0.9, fill=GOLD_500)
    txt(sl, str(i+1), 0.62, y+0.2, 0.45, 0.5,
        size=17, bold=True, color=GOLD_500, align=PP_ALIGN.CENTER)
    txt(sl, val,  1.22, y+0.17, 3.3,  0.55, size=16, bold=True, color=col)
    txt(sl, rest, 4.6,  y+0.2,  7.95, 0.52, size=12, color=GRAY_600)


# ════════════════════════════════════════════════════════
# SLIDE 9 — CIERRE
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg_image(sl, 78)

gradient_title(sl, "LEGENDARYCLASS", 0.5, 0.65, 12.33, 1.55, size=62)

txt(sl, "⚔️  Tu aventura educativa te espera  ⚔️",
    0.5, 2.25, 12.33, 0.5, size=13, color=GRAY_600,
    align=PP_ALIGN.CENTER, italic=True)

gold_line(sl, 3.2, 2.88, 6.93)
card_white(sl, 3.2, 2.98, 6.93, 3.2)

txt(sl, "¡Gracias!", 3.2, 3.1, 6.93, 0.85,
    size=36, bold=True, color=GOLD_600, align=PP_ALIGN.CENTER)
gold_line(sl, 3.6, 4.05, 6.13)
txt(sl, f"{ALUMNO}  ·  {CURSO}  ·  {ANIO}",
    3.2, 4.15, 6.93, 0.38, size=10.5, color=GRAY_600, align=PP_ALIGN.CENTER)
txt(sl, f"{INSTITUCION}  ·  {CARRERA}",
    3.2, 4.55, 6.93, 0.38, size=10.5, color=GRAY_400, align=PP_ALIGN.CENTER)

footer(sl)

# ── Guardar ──────────────────────────────────────────────
OUT = "C:/Proyectos/LegendaryClass/lab03/Presentacion_Caso3_LegendaryClass_v5.pptx"
prs.save(OUT)
print("[OK] PPT guardada:", OUT)
